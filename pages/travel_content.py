"""
✈ AI Travel Banner Studio — Production Grade
=============================================
Outputs TWO things per package:
  1. Direct-post PNG/JPEG (Pillow, 2× supersampled) — pixel-perfect for FB/Instagram
  2. Editable PPTX (python-pptx, single correct slide size) — import to Canva

BUGS FIXED vs uploaded version:
  ✅ PPTX: ONE presentation per format (not mixed sizes — that corrupts the file)
  ✅ PPTX: correct alpha transparency via a:srgbClr/a:alpha XML (not broken _set_shape_alpha)
  ✅ PPTX: no font.color.brightness (not supported in python-pptx)
  ✅ Photos: read once, validated with Image.verify(), cached in session_state
  ✅ QR: generated once, cached, non-blocking (network failure = skip, not crash)
  ✅ Text: Pillow 2× supersampling + stroke_width = antialiased production quality
  ✅ Layout: auto-brightness detection → text placed on darker side

DESIGN: MakeMyTrip / Yatra reference quality
  - Large hero photo full width
  - Bold gold headline with stroke
  - Feature pill grid (STAYS / MEALS / TRANSPORT / GUIDE)
  - Itinerary timeline (city dots)
  - Highlighted price box (special offer style)
  - Footer: social + contact + cert bar

FREE LLMs: Groq (console.groq.com) → Gemini (aistudio.google.com) fallback
"""

import streamlit as st
import sys, os, io, json, re, zipfile, base64
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance, ImageChops, ImageOps

_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)

# ─────────────────────────────────────────────────────────────────────────────
# CONSTANTS
# ─────────────────────────────────────────────────────────────────────────────
SSAA = 2   # 2× supersampling — free antialiasing via LANCZOS downscale

PLATFORMS = {
    "Facebook Post (1200×630)":            (1200, 630),
    "Instagram Post 1:1 (1080×1080)":      (1080, 1080),
    "Instagram / Reels Story (1080×1920)": (1080, 1920),
    "WhatsApp Status (1080×1920)":         (1080, 1920),
    "YouTube Thumbnail (1280×720)":        (1280, 720),
    "Twitter/X Post (1200×675)":           (1200, 675),
}

PPTX_SIZES = {
    "A4 Portrait (8.27×11.69″)":  (8.27, 11.69),
    "Instagram Post (10.8×10.8″)": (10.8, 10.8),
    "Story 9:16 (5.4×9.6″)":       (5.4,   9.6),
    "Facebook (12×6.3″)":          (12.0,  6.3),
    "YouTube Thumbnail (12.8×7.2″)":(12.8,  7.2),
}

# ─────────────────────────────────────────────────────────────────────────────
# DESIGN THEMES
# ─────────────────────────────────────────────────────────────────────────────
THEMES = {
    "🏅 Classic Navy Gold": {
        "primary":    (26,  42,  94),   # deep navy
        "primary2":   (36,  51, 112),   # mid navy
        "accent":     (201,168,  76),   # gold
        "accent2":    (245,216, 120),   # light gold
        "accent_dk":  (139,105,  20),   # dark gold
        "dark":       (10,  14,  40),   # near black navy
        "light":      (255,248, 220),   # cream white
        "pill_colors": [(41,128,185),(39,174,96),(230,126,34),(142,68,173),(192,57,43)],
    },
    "🌿 Emerald Luxury": {
        "primary":    (10,  61,  46),
        "primary2":   (13,  85,  64),
        "accent":     (212,168,  67),
        "accent2":    (245,200,  66),
        "accent_dk":  (140,105,  20),
        "dark":       (4,   40,  30),
        "light":      (240,255, 245),
        "pill_colors": [(13,85,64),(39,174,96),(230,126,34),(0,150,136),(142,68,173)],
    },
    "🌊 Ocean Blue": {
        "primary":    (0,   78, 124),
        "primary2":   (0,  107, 159),
        "accent":     (0,  201, 200),
        "accent2":    (64, 224, 208),
        "accent_dk":  (0,  140, 140),
        "dark":       (0,   48,  85),
        "light":      (240,252, 255),
        "pill_colors": [(0,107,159),(0,201,200),(230,126,34),(39,174,96),(142,68,173)],
    },
    "🌙 Midnight Premium": {
        "primary":    (10,  10,  42),
        "primary2":   (26,  26,  78),
        "accent":     (212,175,  55),
        "accent2":    (255,215,   0),
        "accent_dk":  (160,120,  10),
        "dark":       (0,    0,  16),
        "light":      (255,255, 245),
        "pill_colors": [(26,26,78),(212,175,55),(192,57,43),(39,174,96),(142,68,173)],
    },
    "🌺 Coral Vibrant": {
        "primary":    (176,  48,  32),
        "primary2":   (212,  64,  48),
        "accent":     (245, 166,  35),
        "accent2":    (255,208,  80),
        "accent_dk":  (180, 100,   0),
        "dark":       (120,  24,  16),
        "light":      (255,248, 240),
        "pill_colors": [(176,48,32),(212,64,48),(245,166,35),(39,174,96),(0,107,159)],
    },
}

# ─────────────────────────────────────────────────────────────────────────────
# FONT CACHE
# ─────────────────────────────────────────────────────────────────────────────
_FC: dict = {}
def F(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    k = (size, bold)
    if k in _FC: return _FC[k]
    paths = [
        f"/usr/share/fonts/truetype/dejavu/DejaVuSans{'-Bold' if bold else ''}.ttf",
        f"/usr/share/fonts/truetype/liberation/LiberationSans{'-Bold' if bold else '-Regular'}.ttf",
        f"/usr/share/fonts/truetype/ubuntu/Ubuntu-{'B' if bold else 'R'}.ttf",
    ]
    for p in paths:
        if os.path.exists(p):
            try:
                f = ImageFont.truetype(p, size)
                _FC[k] = f
                return f
            except Exception:
                pass
    return ImageFont.load_default()

# ─────────────────────────────────────────────────────────────────────────────
# KEY HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _get(sk, sec, ev=""):
    v = st.session_state.get(sk,"").strip()
    if v: return v
    try:
        v = str(st.secrets.get(sec,"")).strip()
        if v: return v
    except Exception:
        pass
    return os.getenv(ev,"").strip()

def _groq_key():   return _get("groq_key",   "GROQ_API_KEY",   "GROQ_API_KEY")
def _gemini_key(): return _get("gemini_key", "GEMINI_API_KEY", "GEMINI_API_KEY")
def _llm_ok():     return bool(_groq_key() or _gemini_key())

# ─────────────────────────────────────────────────────────────────────────────
# LLM
# ─────────────────────────────────────────────────────────────────────────────
def _groq(sys_p, usr_p, tokens=1200):
    key = _groq_key()
    if not key: return ""
    try:
        from groq import Groq
        r = Groq(api_key=key).chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role":"system","content":sys_p},
                      {"role":"user","content":usr_p}],
            max_tokens=tokens, temperature=0.75,
            response_format={"type":"json_object"},
        )
        return r.choices[0].message.content.strip()
    except Exception as e:
        st.toast(f"Groq: {e}", icon="⚠️")
        return ""

def _gemini(combined, tokens=1200):
    key = _gemini_key()
    if not key: return ""
    url = "https://generativelanguage.googleapis.com/v1beta/models/gemini-2.5-flash:generateContent"
    body = {"contents":[{"parts":[{"text":combined}]}],
            "generationConfig":{"maxOutputTokens":tokens,"temperature":0.75,
                                "responseMimeType":"application/json"}}
    try:
        import requests as _req
        r = _req.post(url, params={"key":key}, json=body, timeout=35)
        r.raise_for_status()
        return r.json()["candidates"][0]["content"]["parts"][0]["text"].strip()
    except Exception as e:
        st.toast(f"Gemini: {e}", icon="⚠️")
        return ""

def _parse(raw):
    try:
        clean = re.sub(r"```(?:json)?|```","",raw).strip()
        s = clean.find("{")
        return json.loads(clean[s:]) if s >= 0 else json.loads(clean)
    except Exception:
        return {}

def ai_generate(free_text: str, n_photos: int) -> dict:
    sys_p = """You are a travel marketing AI for premium Indian travel agencies.
From a one-line package description generate complete banner content.
Return ONLY valid JSON with keys:
  package_name (short, e.g. "Golden Rajasthan 7D/6N"),
  destination, headline (ALL CAPS, 6-9 words, with 1 emoji),
  subheadline (max 12 words, use · separators),
  duration, price (e.g. "₹28,777"), price_label ("SPECIAL OFFER"),
  price_note (e.g. "per person twin sharing | Breakfast & Dinner"),
  validity (e.g. "Valid till 30 Sep 2026 | T&C Apply"),
  cta ("BOOK NOW" or "ENQUIRE TODAY"),
  pill_features (array of exactly 4: short labels like "3-Star Hotels", "Guided Tours", "All Meals", "Airport Transfers"),
  highlights (array of 5, max 4 words each),
  inclusions (array of 5 key items),
  itinerary (array of {city, nights, hotel} for each stop, max 4),
  day_plan (array of {day, title, details} for each day, max 5),
  hashtags (string, 8 tags),
  instagram_caption (3 lines with emojis + 5 hashtags),
  facebook_caption (2-3 sentences),
  whatsapp_status (2 lines, emoji-rich),
  youtube_title (SEO, max 60 chars),
  reel_script (3 punchy sentences for 30-sec voiceover),
  theme (one of: Classic Navy Gold/Emerald Luxury/Ocean Blue/Midnight Premium/Coral Vibrant)
"""
    usr_p = f'Package: "{free_text}"\nPhotos: {n_photos}\nGenerate:'
    raw = _groq(sys_p, usr_p) or _gemini(f"{sys_p}\n\nReturn ONLY valid JSON.\n\n{usr_p}")
    data = _parse(raw) if raw else {}
    dest = (free_text.split()[0].title() if free_text else "India")
    defaults = {
        "package_name": f"{dest} Package",
        "destination": dest,
        "headline": f"DISCOVER THE MAGIC OF {dest.upper()} ✨",
        "subheadline": f"Heritage · Culture · Adventure",
        "duration": "7 Days / 6 Nights",
        "price": "₹24,999",
        "price_label": "SPECIAL OFFER",
        "price_note": "per person twin sharing",
        "validity": "Limited Seats | T&C Apply",
        "cta": "BOOK NOW",
        "pill_features": ["3-Star Hotels","All Meals","Guided Tours","Airport Transfers"],
        "highlights": ["Iconic Landmarks","Scenic Views","Cultural Experiences","Local Cuisine","Heritage Sites"],
        "inclusions": ["Airport Transfers","Hotel Stay","Daily Breakfast","Expert Guide","All Sightseeing"],
        "itinerary": [{"city":dest,"nights":3,"hotel":"Premium Hotel / Similar"}],
        "day_plan": [
            {"day":"Day 1","title":"Arrival","details":"Arrive, check in, leisure."},
            {"day":"Day 2","title":"Sightseeing","details":"City tour, major attractions."},
            {"day":"Day 3","title":"Excursion","details":"Scenic drive, photo stops."},
        ],
        "hashtags": f"#{dest.replace(' ','')} #Travel #Holiday #TourPackage #India",
        "instagram_caption": f"✈️ {dest} awaits!\n🌟 Book your dream trip now!\n💫 Limited seats! #{dest.replace(' ','')} #Travel",
        "facebook_caption": f"Explore {dest} with our exclusive package! Contact us to book.",
        "whatsapp_status": f"✈️ {dest} Package!\n📞 Book Now!",
        "youtube_title": f"{dest} Tour Package 2025 | Best Deals",
        "reel_script": f"{dest} is calling! From stunning views to rich culture, we cover it all. Book today!",
        "theme": "Classic Navy Gold",
    }
    for k,v in defaults.items():
        data.setdefault(k,v)
    if data.get("theme","") not in THEMES:
        data["theme"] = "Classic Navy Gold"
    return data

# ─────────────────────────────────────────────────────────────────────────────
# IMAGE HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _cover(img: Image.Image, w: int, h: int) -> Image.Image:
    img = img.convert("RGB")
    r = max(w/img.width, h/img.height)
    nw, nh = int(img.width*r), int(img.height*r)
    img = img.resize((nw,nh), Image.LANCZOS)
    l,t = (nw-w)//2, (nh-h)//2
    return img.crop((l,t,l+w,t+h))

def _enhance(img: Image.Image) -> Image.Image:
    img = ImageEnhance.Brightness(img).enhance(1.05)
    img = ImageEnhance.Contrast(img).enhance(1.10)
    img = ImageEnhance.Color(img).enhance(1.15)
    img = ImageEnhance.Sharpness(img).enhance(1.08)
    return img

def detect_text_side(img_bytes: bytes) -> str:
    """Return 'left' or 'right' based on which half is darker (text goes on darker side)."""
    try:
        img = Image.open(io.BytesIO(img_bytes)).convert("L").resize((400,200))
        w,h = img.size
        left_avg  = sum(img.crop((0,0,w//2,h)).getdata()) / (w//2*h)
        right_avg = sum(img.crop((w//2,0,w,h)).getdata()) / (w//2*h)
        return "left" if left_avg < right_avg else "right"
    except Exception:
        return "left"

def _fetch_qr_cached(url: str) -> bytes | None:
    """Fetch QR code, cache in session_state, never crash."""
    cache_key = f"qr_{url}"
    if st.session_state.get(cache_key):
        return st.session_state[cache_key]
    try:
        import requests as _req
        from urllib.parse import quote
        api = f"https://api.qrserver.com/v1/create-qr-code/?size=200x200&data={quote(url)}&bgcolor=ffffff&color=000000&margin=4"
        r = _req.get(api, timeout=6)
        if r.status_code == 200:
            st.session_state[cache_key] = r.content
            return r.content
    except Exception:
        pass
    return None

def load_photos(raw_bytes_list: list) -> list:
    out = []
    for b in raw_bytes_list:
        try:
            img = Image.open(io.BytesIO(bytes(b))).convert("RGB")
            out.append(_enhance(img))
        except Exception:
            pass
    return out

# ─────────────────────────────────────────────────────────────────────────────
# PILLOW DRAWING HELPERS
# ─────────────────────────────────────────────────────────────────────────────
def _tw(draw, xy, text, font, fill, sw=0, sf=(0,0,0)):
    if sw: draw.text(xy, text, font=font, fill=fill, stroke_width=sw, stroke_fill=sf)
    else:  draw.text(xy, text, font=font, fill=fill)

def _mw(draw, xy, text, font, fill, spacing=10, sw=0, sf=(0,0,0)):
    if sw: draw.multiline_text(xy, text, font=font, fill=fill, spacing=spacing, stroke_width=sw, stroke_fill=sf)
    else:  draw.multiline_text(xy, text, font=font, fill=fill, spacing=spacing)

def cx(draw, text, font, y, W, fill, sw=0, sf=(0,0,0)):
    bb = draw.textbbox((0,0),text,font=font)
    tw,th = bb[2]-bb[0], bb[3]-bb[1]
    _tw(draw, ((W-tw)//2, y), text, font, fill, sw, sf)
    return th

def _wrap(text, font, max_w, draw):
    words=text.split(); lines=[]; line=""
    for word in words:
        test=(line+" "+word).strip()
        if draw.textbbox((0,0),test,font=font)[2]<=max_w: line=test
        else:
            if line: lines.append(line)
            line=word
    if line: lines.append(line)
    return "\n".join(lines)

def _pill(draw, x, y, text, font, bg, fg, px=20, py=10, radius=None):
    bb = draw.textbbox((0,0),text,font=font)
    tw,th = bb[2]-bb[0], bb[3]-bb[1]
    r = radius or (th+py)//2
    draw.rounded_rectangle([x,y,x+tw+px*2,y+th+py*2], radius=r, fill=bg)
    _tw(draw, (x+px,y+py), text, font, fg)
    return x+tw+px*2

def grad_v(canvas, x1, y1, x2, y2, c1, c2, alpha_start=None, alpha_end=None):
    """Draw vertical gradient rectangle."""
    draw = ImageDraw.Draw(canvas)
    h = y2-y1
    for y in range(h):
        t = y/h
        r = int(c1[0]+(c2[0]-c1[0])*t)
        g = int(c1[1]+(c2[1]-c1[1])*t)
        b = int(c1[2]+(c2[2]-c1[2])*t)
        if alpha_start is not None and alpha_end is not None:
            a = int(alpha_start+(alpha_end-alpha_start)*t)
            draw.line([(x1,y1+y),(x2,y1+y)], fill=(r,g,b,a))
        else:
            draw.line([(x1,y1+y),(x2,y1+y)], fill=(r,g,b))

def alpha_rect(canvas: Image.Image, x1, y1, x2, y2, color_rgb, alpha: int):
    """Draw semi-transparent rectangle."""
    w, h = x2-x1, y2-y1
    ov = Image.new("RGBA",(w,h), color_rgb+(alpha,))
    canvas.paste(ov, (x1,y1), ov)

# ─────────────────────────────────────────────────────────────────────────────
# MULTI-PHOTO COLLAGE (auto layout)
# ─────────────────────────────────────────────────────────────────────────────
def _collage(canvas: Image.Image, photos: list, x, y, w, h, gap=6):
    n = len(photos)
    if n == 0: return
    if n == 1:
        canvas.paste(_cover(photos[0],w,h), (x,y))
        return
    if n == 2:
        pw=(w-gap)//2
        canvas.paste(_cover(photos[0],pw,h), (x,y))
        canvas.paste(_cover(photos[1],pw,h), (x+pw+gap,y))
        return
    if n == 3:
        lw=int(w*0.60); rw=w-lw-gap; rh=(h-gap)//2
        canvas.paste(_cover(photos[0],lw,h), (x,y))
        canvas.paste(_cover(photos[1],rw,rh), (x+lw+gap,y))
        canvas.paste(_cover(photos[2],rw,rh), (x+lw+gap,y+rh+gap))
        return
    if n == 4:
        lw=int(w*0.55); rw=w-lw-gap; rh=(h-gap)//2
        bw=(rw-gap)//2
        canvas.paste(_cover(photos[0],lw,h), (x,y))
        canvas.paste(_cover(photos[1],rw,rh), (x+lw+gap,y))
        canvas.paste(_cover(photos[2],bw,rh), (x+lw+gap,y+rh+gap))
        canvas.paste(_cover(photos[3],bw,rh), (x+lw+gap+bw+gap,y+rh+gap))
        return
    # 5+: mosaic
    tw=(w-gap*2)//3; th=(h-gap)//2
    canvas.paste(_cover(photos[0],tw*2+gap,th), (x,y))
    canvas.paste(_cover(photos[1],tw,th), (x+tw*2+gap*2,y))
    bw=(w-gap*2)//3
    for i in range(min(3, n-2)):
        canvas.paste(_cover(photos[2+i],bw,th), (x+i*(bw+gap),y+th+gap))

# ─────────────────────────────────────────────────────────────────────────────
# MASTER PILLOW BANNER RENDERER
# ─────────────────────────────────────────────────────────────────────────────
def render_banner(
    photos: list,         # list of PIL Images
    W: int, H: int,
    theme_name: str,
    content: dict,
    brand: dict,
    logo_bytes: bytes | None = None,
    qr_bytes: bytes | None = None,
) -> Image.Image:

    T = THEMES.get(theme_name, THEMES["🏅 Classic Navy Gold"])
    pr  = T["primary"]
    pr2 = T["primary2"]
    ac  = T["accent"]
    ac2 = T["accent2"]
    adk = T["accent_dk"]
    dk  = T["dark"]
    lt  = T["light"]
    pill_cols = T["pill_colors"]

    # Work at 2× for antialiasing
    w, h = W*SSAA, H*SSAA
    sc = min(w, h) / 1080    # scale factor

    canvas = Image.new("RGBA",(w,h),(255,255,255,255))
    draw   = ImageDraw.Draw(canvas)

    # ── DETERMINE LAYOUT: landscape vs portrait ────────────────────────────
    is_landscape = W > H
    is_square    = abs(W-H) < 100

    # ── HERO SECTION ──────────────────────────────────────────────────────
    if is_landscape:
        hero_h = int(h * 0.65)
    elif is_square:
        hero_h = int(h * 0.50)
    else:
        hero_h = int(h * 0.48)

    # Photo background
    if photos:
        hero_bg = _cover(photos[0], w, hero_h)
        canvas.paste(hero_bg.convert("RGBA"), (0,0))
    else:
        # Gradient placeholder
        img_bg = Image.new("RGBA",(w,hero_h))
        grad_v(img_bg, 0, 0, w, hero_h, pr2, dk)
        canvas.paste(img_bg, (0,0))

    # Detect text side and apply directional overlay
    text_side = "left"
    if photos:
        import io as _io
        buf = _io.BytesIO()
        photos[0].save(buf, format="JPEG", quality=85)
        text_side = detect_text_side(buf.getvalue())

    # Gradient overlay — heavier on text side
    ov = Image.new("RGBA",(w,hero_h),(0,0,0,0))
    ovd = ImageDraw.Draw(ov)
    for y in range(hero_h):
        t_frac = (y/hero_h)**1.4
        a_base = int(40+(180-40)*t_frac)
        ovd.line([(0,y),(w,y)], fill=(dk[0],dk[1],dk[2],a_base))
    canvas.alpha_composite(ov)

    # Side panel for text (darker panel on text side)
    if text_side=="left":
        panel_w = int(w*0.58)
        side_ov = Image.new("RGBA",(panel_w,hero_h),(dk[0],dk[1],dk[2],0))
        side_d  = ImageDraw.Draw(side_ov)
        for x_px in range(panel_w):
            t_x = 1-(x_px/panel_w)
            a = int(160*t_x**0.8)
            side_d.line([(x_px,0),(x_px,hero_h)], fill=(dk[0],dk[1],dk[2],a))
        canvas.alpha_composite(side_ov, (0,0))
    else:
        panel_w = int(w*0.58)
        panel_x = w-panel_w
        side_ov = Image.new("RGBA",(panel_w,hero_h),(dk[0],dk[1],dk[2],0))
        side_d  = ImageDraw.Draw(side_ov)
        for x_px in range(panel_w):
            t_x = x_px/panel_w
            a = int(160*t_x**0.8)
            side_d.line([(x_px,0),(x_px,hero_h)], fill=(dk[0],dk[1],dk[2],a))
        canvas.alpha_composite(side_ov, (panel_x,0))

    # Bottom gradient blend into content
    blend_h = int(hero_h * 0.35)
    blend_ov = Image.new("RGBA",(w,blend_h),(dk[0],dk[1],dk[2],0))
    blend_d  = ImageDraw.Draw(blend_ov)
    for y in range(blend_h):
        a = int(220*(y/blend_h)**1.5)
        blend_d.line([(0,y),(w,y)], fill=(dk[0],dk[1],dk[2],a))
    canvas.alpha_composite(blend_ov, (0,hero_h-blend_h))

    draw = ImageDraw.Draw(canvas)

    # Gold accent top bar
    draw.rectangle([0,0,w,int(8*sc)], fill=ac+(255,))

    # ── LOGO ──────────────────────────────────────────────────────────────
    if logo_bytes:
        try:
            logo = Image.open(io.BytesIO(logo_bytes)).convert("RGBA")
            max_dim = int(130*sc)
            r_logo = min(max_dim/logo.width, max_dim/logo.height)
            nw,nh = int(logo.width*r_logo), int(logo.height*r_logo)
            logo = logo.resize((nw,nh), Image.LANCZOS)
            if text_side=="left":
                canvas.paste(logo, (int(50*sc), int(30*sc)), logo)
            else:
                canvas.paste(logo, (w-nw-int(50*sc), int(30*sc)), logo)
        except Exception:
            pass

    # ── HERO TEXT ─────────────────────────────────────────────────────────
    margin = int(52*sc)
    if text_side=="right":
        tx_start = w - int(w*0.56)
    else:
        tx_start = margin

    text_w = int(w*0.54)

    # Package name pill
    pkg = content.get("package_name","")
    if pkg:
        fPkg = F(int(20*sc), bold=True)
        pill_bg = ac+(210,)
        draw.rounded_rectangle(
            [tx_start, int(30*sc),
             tx_start + draw.textbbox((0,0),f"  ✈  {pkg.upper()}  ",font=fPkg)[2]+int(36*sc),
             int(30*sc)+draw.textbbox((0,0),f"  ✈  {pkg.upper()}  ",font=fPkg)[3]+int(18*sc)],
            radius=int(20*sc), fill=pill_bg)
        _tw(draw,(tx_start+int(18*sc),int(36*sc)),f"  ✈  {pkg.upper()}  ",fPkg,dk+(255,))

    # Gold accent line
    line_y = int(90*sc)
    draw.rectangle([tx_start, line_y, tx_start+int(80*sc), line_y+int(5*sc)], fill=ac+(255,))
    line_y += int(20*sc)

    # Headline
    fHL = F(int(74*sc), bold=True)
    headline = content.get("headline","")
    wrapped_hl = _wrap(headline, fHL, text_w, draw)
    _mw(draw, (tx_start, line_y), wrapped_hl, fHL, lt+(255,), spacing=int(10*sc), sw=3, sf=(0,0,0,200))
    bb = draw.multiline_textbbox((tx_start,line_y), wrapped_hl, font=fHL)
    line_y = bb[3] + int(14*sc)

    # Subheadline
    sub = content.get("subheadline","")
    if sub:
        fSub = F(int(36*sc))
        wrapped_sub = _wrap(sub, fSub, text_w, draw)
        _mw(draw,(tx_start,line_y),wrapped_sub,fSub,ac+(255,),spacing=int(6*sc),sw=1,sf=(0,0,0,150))
        bb = draw.multiline_textbbox((tx_start,line_y),wrapped_sub,font=fSub)
        line_y = bb[3] + int(18*sc)

    # Duration badge
    dur = content.get("duration","")
    if dur:
        fDur = F(int(26*sc), bold=True)
        dur_text = f"⏱  {dur}"
        bb = draw.textbbox((0,0),dur_text,font=fDur)
        tw,th = bb[2]-bb[0], bb[3]-bb[1]
        draw.rounded_rectangle([tx_start, line_y, tx_start+tw+int(32*sc), line_y+th+int(16*sc)],
                                 radius=int(12*sc), fill=(255,255,255,55))
        _tw(draw,(tx_start+int(16*sc),line_y+int(8*sc)),dur_text,fDur,lt+(255,),sw=1,sf=(0,0,0,120))

    # ── CONTENT SECTION ────────────────────────────────────────────────────
    cy = hero_h

    # Background for content
    draw.rectangle([0,cy,w,h], fill=lt+(255,))

    # ── PILL FEATURES BAR ──────────────────────────────────────────────────
    pf_h = int(88*sc)
    draw.rectangle([0,cy,w,cy+pf_h], fill=pr+(255,))
    draw.rectangle([0,cy,w,cy+int(4*sc)], fill=ac+(255,))

    features = content.get("pill_features",["Hotels","Meals","Transport","Guide"])[:4]
    fFeat = F(int(22*sc), bold=True)
    fFSub = F(int(18*sc))
    icons = ["🏨","🍽️","🚌","🗺️","✈️","🏖️"]
    fw = (w - int(32*sc)) // len(features)
    for i, feat in enumerate(features):
        fx = int(16*sc) + i*fw
        fy = cy + int(14*sc)
        col = pill_cols[i % len(pill_cols)]
        col_lt = tuple(min(255,c+50) for c in col)
        # Circle icon
        cr = int(22*sc)
        draw.ellipse([fx+int(fw//2)-cr, fy, fx+int(fw//2)+cr, fy+cr*2],
                      fill=col_lt+(200,), outline=ac+(180,), width=int(2*sc))
        icon = icons[i % len(icons)]
        ibb = draw.textbbox((0,0),icon,font=fFSub)
        iw,ih = ibb[2]-ibb[0], ibb[3]-ibb[1]
        draw.text((fx+int(fw//2)-iw//2, fy+cr-ih//2), icon, font=fFSub, fill=lt+(255,))
        # Feature text
        fbb = draw.textbbox((0,0),feat,font=fFeat)
        ftw = fbb[2]-fbb[0]
        ftx = fx + (fw-ftw)//2
        _tw(draw,(ftx,fy+cr*2+int(8*sc)),feat,fFeat,lt+(255,),sw=1,sf=(0,0,0,120))
    cy += pf_h + int(16*sc)

    # ── PHOTO COLLAGE (extra photos) ──────────────────────────────────────
    if len(photos) > 1:
        collage_h = int(240*sc) if not is_landscape else int(180*sc)
        collage_canvas = canvas.crop((0,cy,w,cy+collage_h)).copy()
        extra = photos[1:]
        _collage(collage_canvas, extra, 0, 0, w-int(32*sc), collage_h, gap=int(6*sc))
        # Paste with margin
        canvas.paste(collage_canvas, (int(16*sc),cy))
        # Border
        draw.rectangle([int(16*sc),cy,w-int(16*sc),cy+collage_h],
                         outline=pr+(180,), width=int(3*sc))
        cy += collage_h + int(20*sc)
        draw = ImageDraw.Draw(canvas)

    # ── HIGHLIGHTS ────────────────────────────────────────────────────────
    hls = content.get("highlights",[])
    if hls:
        fHLab = F(int(18*sc), bold=True)
        _tw(draw,(int(16*sc),cy),"HIGHLIGHTS",fHLab,T["primary"]+(200,))
        cy += int(28*sc)
        draw.rectangle([int(16*sc),cy,int(16*sc)+int(50*sc),cy+int(3*sc)], fill=ac+(255,))
        cy += int(16*sc)

        fHL2 = F(int(22*sc))
        col_w = (w - int(48*sc)) // 2
        row_h = int(52*sc)
        for i, hl in enumerate(hls[:6]):
            col = i%2; row = i//2
            bx = int(16*sc) + col*(col_w+int(16*sc))
            by = cy + row*row_h
            # Card
            draw.rounded_rectangle([bx,by,bx+col_w,by+row_h-int(6*sc)],
                                     radius=int(8*sc), fill=(255,255,255,255),
                                     outline=pr+(40,), width=int(2*sc))
            # Accent bar left
            draw.rounded_rectangle([bx,by,bx+int(6*sc),by+row_h-int(6*sc)],
                                     radius=int(4*sc), fill=ac+(230,))
            _tw(draw,(bx+int(14*sc),by+int(12*sc)),f"✦  {hl}",fHL2,pr+(255,),sw=0)
        cy += ((len(hls[:6])+1)//2)*row_h + int(20*sc)

    # ── ITINERARY ─────────────────────────────────────────────────────────
    itin = content.get("itinerary",[])
    if itin:
        fILab = F(int(18*sc), bold=True)
        _tw(draw,(int(16*sc),cy),"ITINERARY",fILab,T["primary"]+(200,))
        cy += int(28*sc)
        draw.rectangle([int(16*sc),cy,int(16*sc)+int(50*sc),cy+int(3*sc)], fill=ac+(255,))
        cy += int(16*sc)

        fCity  = F(int(22*sc), bold=True)
        fHotel = F(int(18*sc))
        fNight = F(int(18*sc), bold=True)
        dot_x  = int(48*sc)
        line_x = dot_x + int(14*sc)

        for idx,stop in enumerate(itin[:4]):
            city   = stop.get("city","")
            nights = stop.get("nights","")
            hotel  = stop.get("hotel","")
            # Connector line
            if idx < len(itin)-1:
                draw.rectangle([line_x,cy+int(28*sc),line_x+int(4*sc),cy+int(72*sc)],fill=ac+(180,))
            # Dot
            dot_r = int(20*sc)
            draw.ellipse([dot_x-dot_r,cy+int(4*sc),dot_x+dot_r,cy+int(4*sc)+dot_r*2],
                          fill=ac+(255,), outline=adk+(255,), width=int(3*sc))
            bb = draw.textbbox((0,0),str(nights),font=fNight)
            nw2,nh2 = bb[2]-bb[0],bb[3]-bb[1]
            _tw(draw,(dot_x-nw2//2+int(2*sc),cy+int(4*sc)+dot_r-nh2//2),
                str(nights),fNight,dk+(255,))
            # City + hotel
            _tw(draw,(dot_x+dot_r+int(16*sc),cy+int(6*sc)),city.upper(),fCity,pr+(255,),sw=1,sf=(255,255,255,80))
            _tw(draw,(dot_x+dot_r+int(16*sc),cy+int(34*sc)),hotel or "Premium Hotel / Similar",fHotel,(80,80,80,255))
            cy += int(72*sc)
        cy += int(12*sc)

    # ── DEAL CARD (price box) ──────────────────────────────────────────────
    deal_h = int(180*sc)
    # Gold gradient deal card
    deal_img = Image.new("RGBA",(w-int(32*sc),deal_h))
    grad_v(deal_img, 0, 0, w-int(32*sc), deal_h,
           (max(0,ac[0]-30),max(0,ac[1]-30),max(0,ac[2]-30)),
           (min(255,ac2[0]+20),min(255,ac2[1]+20),min(255,ac2[2]+20)))
    canvas.paste(deal_img,(int(16*sc),cy),deal_img)
    draw = ImageDraw.Draw(canvas)

    # Border
    draw.rounded_rectangle([int(16*sc),cy,w-int(16*sc),cy+deal_h],
                             radius=int(14*sc), outline=adk+(255,), width=int(3*sc))

    # Label pill
    label = content.get("price_label","SPECIAL OFFER")
    fLabel = F(int(20*sc), bold=True)
    lbb = draw.textbbox((0,0),label,font=fLabel)
    lw = lbb[2]-lbb[0]; lh = lbb[3]-lbb[1]
    draw.rounded_rectangle([int(40*sc),cy+int(16*sc),int(40*sc)+lw+int(32*sc),cy+int(16*sc)+lh+int(16*sc)],
                             radius=int(10*sc), fill=pr+(230,))
    _tw(draw,(int(56*sc),cy+int(24*sc)),label,fLabel,ac2+(255,),sw=0)

    # Price — BIG
    price = content.get("price","₹24,999")
    fPrice = F(int(64*sc), bold=True)
    _tw(draw,(int(40*sc),cy+int(56*sc)),price,fPrice,pr+(255,),sw=2,sf=adk+(200,))

    # Price note
    note = content.get("price_note","per person")
    fNote = F(int(20*sc))
    _tw(draw,(int(40*sc),cy+int(130*sc)),note,fNote,(60,40,10,255))

    # CTA button (right side of deal card)
    cta_text = f"  {content.get('cta','BOOK NOW')} →  "
    fCTA = F(int(28*sc), bold=True)
    cta_bb = draw.textbbox((0,0),cta_text,font=fCTA)
    cta_w = cta_bb[2]-cta_bb[0]+int(20*sc)
    cta_h2 = cta_bb[3]-cta_bb[1]+int(20*sc)
    cta_x = w - cta_w - int(40*sc)
    cta_y = cy + (deal_h-cta_h2)//2
    draw.rounded_rectangle([cta_x,cta_y,cta_x+cta_w,cta_y+cta_h2],
                             radius=int(14*sc), fill=pr+(240,))
    draw.rounded_rectangle([cta_x+int(3*sc),cta_y+int(3*sc),cta_x+cta_w-int(3*sc),cta_y+cta_h2-int(3*sc)],
                             radius=int(12*sc), outline=ac+(200,), width=int(2*sc))
    _tw(draw,(cta_x+int(10*sc),cta_y+int(10*sc)),cta_text,fCTA,ac2+(255,),sw=1,sf=adk+(180,))

    cy += deal_h + int(16*sc)

    # ── INCLUSIONS BAR ────────────────────────────────────────────────────
    inclusions = content.get("inclusions",[])
    if inclusions and cy + int(100*sc) < h:
        incl_h = int(90*sc)
        draw.rectangle([0,cy,w,cy+incl_h], fill=pr+(255,))
        draw.rectangle([0,cy,w,cy+int(4*sc)], fill=ac+(255,))
        fInc = F(int(19*sc))
        fIncLab = F(int(16*sc), bold=True)
        _tw(draw,(int(20*sc),cy+int(8*sc)),"INCLUSIONS",fIncLab,ac2+(255,))
        inc_y = cy+int(36*sc)
        inc_cols = 3
        inc_cw   = (w-int(40*sc)) // inc_cols
        for i, inc in enumerate(inclusions[:6]):
            ix = int(16*sc) + (i%inc_cols)*inc_cw
            iy = inc_y + (i//inc_cols)*int(28*sc)
            _tw(draw,(ix,iy),f"✓  {inc}",fInc,lt+(220,),sw=0)
        cy += incl_h

    # ── FOOTER ────────────────────────────────────────────────────────────
    remaining = h - cy
    if remaining < int(130*sc):
        # Extend canvas if needed
        extra = int(130*sc) - remaining + int(20*sc)
        new_canvas = Image.new("RGBA",(w,h+extra),(255,255,255,255))
        new_canvas.paste(canvas,(0,0))
        canvas = new_canvas
        h += extra
        draw = ImageDraw.Draw(canvas)
        draw.rectangle([0,cy,w,h],fill=pr+(255,))

    # Social bar
    soc_y = h - int(110*sc)
    draw.rectangle([0,soc_y,w,soc_y+int(50*sc)],fill=(255,255,255,255))
    draw.rectangle([0,soc_y,w,soc_y+int(3*sc)],fill=pr+(255,))
    fSoc = F(int(18*sc))
    fSocT = F(int(18*sc), bold=True)
    cx(draw,"FOLLOW US",fSocT,soc_y+int(8*sc),w,pr+(255,))
    soc_line = "  ·  ".join(filter(None,[
        f"f  {brand.get('facebook','')}" if brand.get('facebook') else "",
        f"@  {brand.get('instagram','')}" if brand.get('instagram') else "",
        f"▶  {brand.get('youtube','')}"  if brand.get('youtube')   else "",
    ]))
    if soc_line:
        cx(draw,soc_line,fSoc,soc_y+int(28*sc),w,pr+(200,))

    # Contact bar
    contact_y = soc_y + int(50*sc)
    draw.rectangle([0,contact_y,w,h-int(32*sc)],fill=pr2+(255,))
    fCt = F(int(20*sc), bold=True)
    fCtSub = F(int(18*sc))
    web   = brand.get("website","")
    phone = brand.get("phone","")
    e1    = brand.get("email1","")
    if web or phone:
        cx(draw,f"🌐 {web}   📞 {phone}",fCt,contact_y+int(10*sc),w,lt+(255,))
    if e1:
        cx(draw,e1,fCtSub,contact_y+int(38*sc),w,ac2+(220,))

    # QR code
    if qr_bytes:
        try:
            qr = Image.open(io.BytesIO(qr_bytes)).convert("RGB")
            qr_size = int(80*sc)
            qr = qr.resize((qr_size,qr_size),Image.LANCZOS)
            canvas.paste(qr, (w-qr_size-int(20*sc), contact_y+int(8*sc)))
        except Exception:
            pass

    # Cert bar
    draw.rectangle([0,h-int(32*sc),w,h],fill=lt+(255,))
    draw.rectangle([0,h-int(32*sc),w,h-int(32*sc)+int(2*sc)],fill=ac+(180,))
    certs = brand.get("certs",["IATA","OTAI","ADTOI","NIMA","ETAA"])
    fCert = F(int(14*sc), bold=True)
    cert_str = "   ·   ".join(certs)
    cx(draw,cert_str,fCert,h-int(26*sc),w,pr+(180,))

    # Copyright
    draw.rectangle([0,h-int(28*sc)+int(2*sc),w,h-int(28*sc)+int(4*sc)],fill=(220,220,220,255))
    # Gold frame border
    draw.rectangle([0,0,w,int(6*sc)],fill=ac+(255,))
    draw.rectangle([0,h-int(6*sc),w,h],fill=ac+(255,))

    # Final downscale → free LANCZOS antialiasing
    final = canvas.convert("RGB").resize((W,H), Image.LANCZOS)
    return final

def to_bytes(img, fmt="PNG", quality=95):
    buf = io.BytesIO()
    img.convert("RGB").save(buf, format=fmt, quality=quality)
    return buf.getvalue()

# ─────────────────────────────────────────────────────────────────────────────
# PPTX BUILDER (FIXED — single correct slide size)
# ─────────────────────────────────────────────────────────────────────────────
def build_pptx(content, theme_name, photos_raw, logo_bytes, pptx_size_name):
    """
    Fixed PPTX builder:
    - ONE presentation = ONE slide size (set BEFORE adding any slides)
    - Correct alpha via a:srgbClr/a:alpha XML
    - No font.color.brightness (not supported)
    - One PPTX per platform size
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    from lxml import etree

    T = THEMES.get(theme_name, THEMES["🏅 Classic Navy Gold"])

    def _rgb(hex_str):
        h = hex_str.lstrip("#")
        if len(h)==6: return RGBColor(int(h[0:2],16),int(h[2:4],16),int(h[4:6],16))
        r,g,b = hex_str if isinstance(hex_str,tuple) else (26,42,94)
        return RGBColor(r,g,b)

    def rgb_to_hex(rgb):
        return "%02x%02x%02x" % (rgb[0],rgb[1],rgb[2])

    def set_alpha(shape, alpha_pct):
        """Set shape fill transparency 0-100 (0=transparent, 100=opaque)."""
        val = int(alpha_pct * 1000)  # OOXML uses 0-100000
        try:
            sp = shape._element.spPr
            solid = sp.find('.//' + qn('a:solidFill'))
            if solid is not None:
                srgb = solid.find(qn('a:srgbClr'))
                if srgb is not None:
                    for old in srgb.findall(qn('a:alpha')):
                        srgb.remove(old)
                    alpha_el = etree.SubElement(srgb, qn('a:alpha'))
                    alpha_el.set('val', str(val))
        except Exception:
            pass

    def no_line(shape):
        try:
            shape.line.fill.background()
            shape.line.width = Pt(0)
        except Exception:
            pass

    def add_rect(slide, x, y, w2, h2, rgb, alpha=100):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                      Inches(x),Inches(y),Inches(w2),Inches(h2))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*rgb)
        no_line(shp)
        if alpha < 100:
            set_alpha(shp, alpha)
        return shp

    def add_rr(slide, x, y, w2, h2, rgb, alpha=100):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                      Inches(x),Inches(y),Inches(w2),Inches(h2))
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(*rgb)
        no_line(shp)
        if alpha < 100:
            set_alpha(shp, alpha)
        return shp

    def add_txt(slide, text, x, y, w2, h2, size=12, bold=False,
                rgb=(0,0,0), align=PP_ALIGN.LEFT, font="Aptos"):
        tb = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w2),Inches(h2))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left=tf.margin_right=Inches(0.06)
        tf.margin_top=tf.margin_bottom=Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = RGBColor(*rgb)
        run.font.name  = font
        return tb

    def set_txt_in_shape(shp, text, size=12, bold=False, rgb=(255,255,255), align=PP_ALIGN.CENTER):
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left=tf.margin_right=Inches(0.06)
        tf.margin_top=tf.margin_bottom=Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = RGBColor(*rgb)

    def add_pic(slide, img_bytes, x, y, w2, h2):
        try:
            img = Image.open(io.BytesIO(bytes(img_bytes))).convert("RGB")
            img = ImageOps.fit(img,(int(w2*150),int(h2*150)),Image.LANCZOS)
            buf = io.BytesIO(); img.save(buf,format="JPEG",quality=88)
            buf.seek(0)
            slide.shapes.add_picture(buf,Inches(x),Inches(y),Inches(w2),Inches(h2))
        except Exception:
            pass

    pr  = T["primary"]
    pr2 = T["primary2"]
    ac  = T["accent"]
    ac2 = T["accent2"]
    dk  = T["dark"]
    lt  = T["light"]

    # ── PPTX — set size BEFORE adding slides ──────────────────────────────
    prs = Presentation()
    SW, SH = PPTX_SIZES.get(pptx_size_name, (8.27,11.69))
    prs.slide_width  = Inches(SW)
    prs.slide_height = Inches(SH)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    W = SW

    # White background
    add_rect(slide,0,0,W,SH,(255,255,255))

    # Hero photo
    HERO_H = SH*0.42
    if photos_raw:
        add_pic(slide, photos_raw[0], 0, 0, W, HERO_H)

    # Dark overlay on hero
    ov = add_rect(slide, 0, 0, W, HERO_H, dk, alpha=65)

    # Side panel overlay
    if photos_raw:
        buf_tmp = io.BytesIO()
        Image.open(io.BytesIO(bytes(photos_raw[0]))).save(buf_tmp,format="JPEG",quality=70)
        ts = detect_text_side(buf_tmp.getvalue())
    else:
        ts = "left"

    if ts=="left":
        sp_ov = add_rect(slide,0,0,W*0.6,HERO_H,dk,alpha=55)
    else:
        sp_ov = add_rect(slide,W*0.4,0,W*0.6,HERO_H,dk,alpha=55)

    # Gold top bar
    add_rect(slide,0,0,W,0.06,ac)

    # Logo
    if logo_bytes:
        try:
            x_logo = 0.25 if ts=="left" else W-1.45
            add_pic(slide,logo_bytes,x_logo,0.18,1.2,0.6)
        except Exception:
            pass

    tx_x = 0.3 if ts=="left" else W-4.0

    # Package name badge
    pkg = content.get("package_name","")
    if pkg:
        badge = add_rr(slide,tx_x,0.28,3.2,0.38,ac)
        set_txt_in_shape(badge,f"✈  {pkg.upper()}",size=10,bold=True,rgb=dk,align=PP_ALIGN.LEFT)

    # Gold accent line
    add_rect(slide,tx_x,0.74,0.6,0.04,ac)

    # Headline
    hl = content.get("headline","")
    hl_size = max(16, int(32*(20/max(len(hl),1))**0.5))
    add_txt(slide,hl,tx_x,0.82,3.6,1.0,size=min(hl_size,28),bold=True,
            rgb=lt,font="Aptos Display",
            align=PP_ALIGN.LEFT if ts=="left" else PP_ALIGN.RIGHT)

    add_txt(slide,content.get("subheadline",""),tx_x,HERO_H-0.7,3.6,0.4,
            size=11,rgb=ac2,
            align=PP_ALIGN.LEFT if ts=="left" else PP_ALIGN.RIGHT)

    # Duration
    dur_box = add_rr(slide,tx_x,HERO_H-0.42,2.2,0.36,(255,255,255),alpha=25)
    set_txt_in_shape(dur_box,f"⏱ {content.get('duration','')}",size=11,bold=True,
                     rgb=lt,align=PP_ALIGN.LEFT)

    CY = HERO_H + 0.18

    # Feature pills row
    features = content.get("pill_features",["Hotels","Meals","Guides","Transfers"])[:4]
    add_rect(slide,0,CY,W,0.72,pr)
    add_rect(slide,0,CY,W,0.04,ac)
    fw = (W-0.4)/len(features)
    for i,feat in enumerate(features):
        fx = 0.2+i*fw
        col = T["pill_colors"][i%len(T["pill_colors"])]
        pill = add_rr(slide,fx,CY+0.12,fw-0.15,0.48,col)
        set_txt_in_shape(pill,feat,size=10,bold=True,rgb=lt)
    CY += 0.78

    # Collage
    if len(photos_raw)>1:
        ch = min(1.8, SH-CY-3.5)
        n_photos = min(len(photos_raw)-1,4)
        cw=(W-0.4)/n_photos
        for i in range(n_photos):
            add_pic(slide,photos_raw[i+1],0.2+i*cw,CY,cw-0.05,ch)
        CY += ch+0.18

    # Highlights grid
    hls = content.get("highlights",[])
    if hls:
        add_txt(slide,"HIGHLIGHTS",0.2,CY,W-0.4,0.26,size=9,bold=True,rgb=T["primary"]+(200,200,200))
        CY += 0.3
        add_rect(slide,0.2,CY,0.4,0.03,ac)
        CY += 0.12
        col_w=(W-0.5)/2
        for i,hl in enumerate(hls[:6]):
            col=i%2; row=i//2
            bx=0.2+col*col_w; by=CY+row*0.45
            card=add_rr(slide,bx,by,col_w-0.1,0.4,(255,255,255))
            bar=add_rect(slide,bx,by,0.05,0.4,ac)
            add_txt(slide,f"✦ {hl}",bx+0.1,by+0.06,col_w-0.25,0.3,
                    size=9,bold=True,rgb=pr)
        CY += ((len(hls[:6])+1)//2)*0.48+0.12

    # Deal card
    deal_h = 1.3
    deal = add_rr(slide,0.2,CY,W-0.4,deal_h,ac)
    lbl = add_rr(slide,0.35,CY+0.12,2.0,0.32,pr)
    set_txt_in_shape(lbl,content.get("price_label","SPECIAL OFFER"),
                     size=9,bold=True,rgb=ac2,align=PP_ALIGN.LEFT)
    add_txt(slide,content.get("price","₹24,999"),0.35,CY+0.46,3.5,0.6,
            size=32,bold=True,rgb=dk,font="Aptos Display")
    add_txt(slide,content.get("price_note",""),0.35,CY+1.05,3.5,0.22,
            size=8,rgb=(60,40,10))
    cta_b = add_rr(slide,W-2.6,CY+0.38,2.2,0.56,dk)
    set_txt_in_shape(cta_b,f"{content.get('cta','BOOK NOW')} →",
                     size=13,bold=True,rgb=ac2)
    CY += deal_h+0.18

    # Inclusions
    inclusions=content.get("inclusions",[])
    if inclusions:
        incl_panel=add_rect(slide,0,CY,W,0.9,pr2)
        add_txt(slide,"INCLUSIONS",0.2,CY+0.06,W-0.4,0.22,size=8,bold=True,rgb=ac2)
        cw2=(W-0.5)/3
        for i,inc in enumerate(inclusions[:6]):
            ix=0.2+(i%3)*cw2; iy=CY+0.32+(i//3)*0.26
            add_txt(slide,f"✓ {inc}",ix,iy,cw2-0.05,0.24,size=8,rgb=lt)
        CY += 0.95

    # Footer
    footer_h = SH-CY
    add_rect(slide,0,CY,W,footer_h,dk)
    add_txt(slide,f"📞 {content.get('phone','')}   🌐 {content.get('website','')}",
            0.2,CY+0.12,W-0.4,0.28,size=9,bold=True,rgb=lt,align=PP_ALIGN.CENTER)

    qr = _fetch_qr_cached(content.get("website","www.example.com"))
    if qr:
        add_pic(slide,qr,W-1.35,CY+0.08,1.1,1.1)

    certs = content.get("_brand",{}).get("certs",["IATA","OTAI","ADTOI","NIMA","ETAA"])
    cert_str = "   ·   ".join(certs)
    add_txt(slide,cert_str,0.2,SH-0.32,W-0.4,0.28,size=7,rgb=(200,200,200),align=PP_ALIGN.CENTER)

    # Gold frame
    try:
        from pptx.util import Pt as _Pt
        frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                        Inches(0.06),Inches(0.06),
                                        Inches(W-0.12),Inches(SH-0.12))
        frame.fill.background()
        frame.line.color.rgb = RGBColor(*ac)
        frame.line.width = _Pt(2.5)
    except Exception:
        pass

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()

# ─────────────────────────────────────────────────────────────────────────────
# PAGE
# ─────────────────────────────────────────────────────────────────────────────
def render():
    st.markdown("""<style>
    .hero-h{font-size:2rem;font-weight:800;
      background:linear-gradient(120deg,#c9a84c,#1a2a5e);
      -webkit-background-clip:text;-webkit-text-fill-color:transparent}
    .hero-s{color:#6b7280;font-size:.88rem;margin-top:2px}
    .bdg{display:inline-block;padding:3px 12px;border-radius:20px;
         font-size:.7rem;font-weight:700;margin-right:6px;margin-bottom:8px}
    .b1{background:#1a2a5e;color:#c9a84c}
    .b2{background:#c9a84c;color:#1a2a5e}
    .b3{background:linear-gradient(135deg,#7c3aed,#db2777);color:#fff}
    .empty{border:1px dashed #374151;border-radius:12px;padding:60px 20px;text-align:center}
    .ei{font-size:2.8rem}.et{color:#6b7280;margin-top:8px;font-size:.85rem}
    </style>""", unsafe_allow_html=True)

    st.markdown('<span class="bdg b1">✈ TRAVEL</span>'
                '<span class="bdg b2">🏆 PRODUCTION</span>'
                '<span class="bdg b3">🤖 AI</span>', unsafe_allow_html=True)
    st.markdown('<div class="hero-h">AI Travel Banner Studio</div>', unsafe_allow_html=True)
    st.markdown('<div class="hero-s">One-prompt → pixel-perfect PNG for FB/Instagram + editable PPTX for Canva</div>',
                unsafe_allow_html=True)
    st.markdown("---")

    # ── BRAND KIT ─────────────────────────────────────────────────────────
    with st.expander("🏷️ Brand Kit — logo · social · contact · certs", expanded=False):
        r1,r2,r3 = st.columns(3)
        with r1:
            st.markdown("**🖼️ Company Logo**")
            lu = st.file_uploader("Logo PNG",type=["png","jpg","jpeg"],key="bk_logo")
            if lu: st.session_state["brand_logo"] = lu.read()
            if st.session_state.get("brand_logo"):
                st.image(st.session_state["brand_logo"],width=110)
                if st.button("✕",key="rm_logo"): del st.session_state["brand_logo"]

        with r2:
            st.markdown("**🏢 Company Details**")
            bn   = st.text_input("Company Name",  value=st.session_state.get("bk_name","7 WONDERS WORLD"),key="_bn")
            bweb = st.text_input("Website",        value=st.session_state.get("bk_web","www.7wwtravels.com"),key="_bweb")
            bph  = st.text_input("Phone",          value=st.session_state.get("bk_ph","+91 97112 81598"),key="_bph")
            be1  = st.text_input("Email",          value=st.session_state.get("bk_e1","info@7wwtravels.com"),key="_be1")
            bcerts=st.text_input("Certs (comma-sep)",value=st.session_state.get("bk_certs","IATA,OTAI,ADTOI,NIMA,ETAA"),key="_bc")

        with r3:
            st.markdown("**📱 Social Handles**")
            bfb = st.text_input("Facebook", value=st.session_state.get("bk_fb","7wwtravels"),key="_bfb")
            big = st.text_input("Instagram",value=st.session_state.get("bk_ig","7ww_travels"),key="_big")
            byt = st.text_input("YouTube",  value=st.session_state.get("bk_yt","@7wwtravels"),key="_byt")
            if st.button("💾 Save Brand Kit",use_container_width=True):
                st.session_state.update(bk_name=bn,bk_web=bweb,bk_ph=bph,bk_e1=be1,
                                         bk_certs=bcerts,bk_fb=bfb,bk_ig=big,bk_yt=byt)
                st.success("Brand kit saved!")

        st.markdown("---")
        st.markdown("##### 🔑 AI Keys — free")
        ka,kb = st.columns(2)
        with ka:
            gq = st.text_input("⚡ Groq",type="password",
                                value=st.session_state.get("groq_key",""),
                                placeholder="gsk_...",key="gq_in",
                                help="console.groq.com → free, fast")
            if gq: st.session_state["groq_key"]=gq.strip()
            st.success("✓ Groq active") if _groq_key() else st.info("console.groq.com")
        with kb:
            gm = st.text_input("🔵 Gemini",type="password",
                                value=st.session_state.get("gemini_key",""),
                                placeholder="AIzaSy...",key="gm_in")
            if gm: st.session_state["gemini_key"]=gm.strip()
            st.success("✓ Gemini active") if _gemini_key() else st.info("aistudio.google.com")

    # Brand dict
    brand = {
        "name":      st.session_state.get("bk_name","7 WONDERS WORLD"),
        "website":   st.session_state.get("bk_web","www.7wwtravels.com"),
        "phone":     st.session_state.get("bk_ph","+91 97112 81598"),
        "email1":    st.session_state.get("bk_e1","info@7wwtravels.com"),
        "facebook":  st.session_state.get("bk_fb","7wwtravels"),
        "instagram": st.session_state.get("bk_ig","7ww_travels"),
        "youtube":   st.session_state.get("bk_yt","@7wwtravels"),
        "certs":     [c.strip() for c in st.session_state.get("bk_certs","IATA,OTAI,ADTOI,NIMA,ETAA").split(",") if c.strip()],
    }
    logo_bytes = st.session_state.get("brand_logo")

    st.markdown("---")
    st.markdown("### ✍️ Describe your package in one line")
    free_text = st.text_area("",height=90,key="free_text",
        placeholder="5N/6D Bhutan, Thimphu 2N + Punakha 1N + Paro 2N, Rs 28777 per person twin sharing, breakfast dinner, Toyota, guide, airport transfers, min 4 guests, valid Sep 2026",
        label_visibility="collapsed")

    st.markdown("**📸 Upload Photos (1-6)**")
    photos_input = st.file_uploader("",type=["jpg","jpeg","png","webp"],
                                     accept_multiple_files=True,key="main_photos",
                                     label_visibility="collapsed")
    # Read + validate ALL photos immediately (prevent BytesIO drain)
    if photos_input:
        names = [p.name for p in photos_input]
        if st.session_state.get("_photo_names") != names:
            fresh = []
            for p in photos_input[:6]:
                try:
                    p.seek(0); b = p.read()
                    if b and len(b)>100:
                        Image.open(io.BytesIO(b)).verify()
                        fresh.append(b)
                except Exception:
                    pass
            st.session_state["_photo_bytes"] = fresh
            st.session_state["_photo_names"] = names

    cached_bytes: list = st.session_state.get("_photo_bytes",[])

    if cached_bytes:
        cols = st.columns(min(len(cached_bytes),6))
        for i,b in enumerate(cached_bytes):
            try:
                img = Image.open(io.BytesIO(b))
                s = 100/img.width
                cols[i].image(to_bytes(img.resize((100,int(img.height*s)),Image.LANCZOS)),
                               use_container_width=True)
            except Exception:
                pass

    cg1,cg2 = st.columns([3,1])
    with cg1:
        big_gen = st.button("🚀 Generate All Content",type="primary",
                             use_container_width=True,
                             disabled=not(free_text.strip() and cached_bytes))
    with cg2:
        if not _llm_ok(): st.warning("Add API key ↑")
        elif not free_text.strip(): st.info("Describe package ↑")
        elif not cached_bytes: st.info("Upload photos ↑")

    if big_gen and free_text.strip() and cached_bytes:
        with st.spinner("🤖 AI generating complete package content…"):
            ai_data = ai_generate(free_text, len(cached_bytes))
        ai_data["_brand"] = brand
        st.session_state.update(ai_data=ai_data, ai_photos=cached_bytes)
        st.success(f"✅ {ai_data.get('headline','')}  ·  Theme: {ai_data.get('theme','')}")

    if not _llm_ok() and big_gen and cached_bytes:
        ai_data = ai_generate(free_text, len(cached_bytes))
        ai_data["_brand"] = brand
        st.session_state.update(ai_data=ai_data, ai_photos=cached_bytes)

    ai  = st.session_state.get("ai_data",{})
    raw = st.session_state.get("ai_photos",[])

    if not ai:
        st.markdown("---")
        st.markdown('<div class="empty"><div class="ei">✍️</div>'
                    '<div class="et">Describe your package above → Generate</div></div>',
                    unsafe_allow_html=True)
        return

    st.markdown("---")

    tab_banner, tab_pptx, tab_copy, tab_bulk = st.tabs([
        "🖼️ Banner (PNG/JPEG)",
        "📐 Editable PPTX (Canva)",
        "📋 AI Copy & Captions",
        "📦 Bulk Export",
    ])

    # ── FINE-TUNE SIDEBAR ─────────────────────────────────────────────────
    with st.expander("🔧 Edit AI Content", expanded=False):
        ec1,ec2 = st.columns(2)
        with ec1:
            ai["headline"]    = st.text_input("Headline",   value=ai.get("headline",""),  key="e_hl")
            ai["subheadline"] = st.text_input("Subheadline",value=ai.get("subheadline",""),key="e_sub")
            ai["duration"]    = st.text_input("Duration",   value=ai.get("duration",""),  key="e_dur")
            ai["price"]       = st.text_input("Price",      value=ai.get("price",""),     key="e_price")
            ai["price_label"] = st.text_input("Price Label",value=ai.get("price_label","SPECIAL OFFER"),key="e_pl")
            ai["price_note"]  = st.text_input("Price Note", value=ai.get("price_note",""),key="e_pn")
        with ec2:
            ai["cta"]         = st.text_input("CTA",        value=ai.get("cta","BOOK NOW"),key="e_cta")
            ai["validity"]    = st.text_input("Valid till",  value=ai.get("validity",""),  key="e_val")
            ai["website"]     = st.text_input("Website",    value=brand.get("website",""),key="e_web")
            ai["phone"]       = st.text_input("Phone",      value=brand.get("phone",""),  key="e_ph")
            hl_raw = st.text_area("Highlights (one per line)",
                                   value="\n".join(ai.get("highlights",[])),height=100,key="e_hls")
            ai["highlights"] = [l.strip() for l in hl_raw.splitlines() if l.strip()]

    # ════════════════════════════════════════════════════════════════════
    # TAB 1 — PNG BANNER
    # ════════════════════════════════════════════════════════════════════
    with tab_banner:
        bl,br = st.columns([1,1],gap="large")
        with bl:
            theme_name = st.selectbox("Theme",list(THEMES.keys()),
                index=list(THEMES.keys()).index(ai.get("theme","🏅 Classic Navy Gold"))
                      if ai.get("theme") in THEMES else 0, key="b_theme")
            platform = st.selectbox("Platform",list(PLATFORMS.keys()),key="b_plat")
            gen_banner = st.button("🎨 Render Banner",type="primary",
                                    use_container_width=True,disabled=not raw)

        with br:
            st.markdown("### 👁️ Preview")
            if gen_banner and raw:
                W,H = PLATFORMS[platform]
                photos_pil = load_photos(raw)

                # Fetch QR (cached, non-blocking)
                qr = _fetch_qr_cached(ai.get("website") or brand.get("website",""))

                with st.spinner(f"Rendering {W}×{H} at 2× quality…"):
                    ai["website"] = ai.get("website") or brand.get("website","")
                    ai["phone"]   = ai.get("phone")   or brand.get("phone","")
                    ai["_brand"]  = brand
                    banner = render_banner(photos_pil,W,H,theme_name,ai,brand,logo_bytes,qr)

                st.session_state.update(
                    b_png=to_bytes(banner,"PNG"),
                    b_jpg=to_bytes(banner,"JPEG",93),
                    b_name=f"{ai.get('package_name','banner')}_{platform[:16]}")

            if st.session_state.get("b_png"):
                st.image(st.session_state["b_png"],use_container_width=True)
                d1,d2 = st.columns(2)
                with d1:
                    st.download_button("📥 Download PNG",
                        data=st.session_state["b_png"],
                        file_name=f"{st.session_state['b_name']}.png",
                        mime="image/png",use_container_width=True)
                with d2:
                    st.download_button("📥 Download JPEG",
                        data=st.session_state["b_jpg"],
                        file_name=f"{st.session_state['b_name']}.jpg",
                        mime="image/jpeg",use_container_width=True)
                st.success("✅ Production-quality banner ready!")
            else:
                st.markdown('<div class="empty"><div class="ei">🖼️</div>'
                            '<div class="et">Select platform → Render Banner</div></div>',
                            unsafe_allow_html=True)

    # ════════════════════════════════════════════════════════════════════
    # TAB 2 — PPTX
    # ════════════════════════════════════════════════════════════════════
    with tab_pptx:
        pl,pr2_ = st.columns([1,1],gap="large")
        with pl:
            st.markdown("### 📐 Editable PPTX for Canva")
            st.info("Import into Canva → every element is a separate editable shape")
            pptx_theme = st.selectbox("Theme",list(THEMES.keys()),
                index=list(THEMES.keys()).index(ai.get("theme","🏅 Classic Navy Gold"))
                      if ai.get("theme") in THEMES else 0, key="p_theme")
            pptx_size  = st.selectbox("Slide Size",list(PPTX_SIZES.keys()),key="p_size")
            gen_pptx   = st.button("📐 Generate PPTX",type="primary",
                                    use_container_width=True,disabled=not raw)

        with pr2_:
            st.markdown("### 📥 Download")
            if gen_pptx and raw:
                ai["website"] = ai.get("website") or brand.get("website","")
                ai["phone"]   = ai.get("phone")   or brand.get("phone","")
                ai["_brand"]  = brand
                with st.spinner("Building PPTX…"):
                    pptx_b = build_pptx(ai,pptx_theme,raw,logo_bytes,pptx_size)
                st.session_state["pptx_b"] = pptx_b

            if st.session_state.get("pptx_b"):
                st.success("✅ PPTX ready!")
                st.download_button("📥 Download PPTX (Canva-ready)",
                    data=st.session_state["pptx_b"],
                    file_name=f"{ai.get('package_name','flyer')}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True)
                st.markdown("""
**How to edit in Canva:**
1. Download the PPTX file
2. Open Canva → *Create design* → *Import file* → upload PPTX
3. Every text box, shape and photo is a separate editable element
4. Change fonts, colors, text → Download as PNG/PDF
                """)
            else:
                st.markdown('<div class="empty"><div class="ei">📐</div>'
                            '<div class="et">Click Generate PPTX</div></div>',
                            unsafe_allow_html=True)

    # ════════════════════════════════════════════════════════════════════
    # TAB 3 — AI COPY
    # ════════════════════════════════════════════════════════════════════
    with tab_copy:
        st.markdown("### 📋 All AI-Generated Copy")
        def _cbox(label,val):
            if not val: return
            st.markdown(f"**{label}**")
            v = "\n".join(f"• {x}" for x in val) if isinstance(val,list) else str(val)
            st.text_area("",value=v,height=max(60,min(180,v.count("\n")*28+60)),
                          key=f"cp_{label}",label_visibility="collapsed")
        c1,c2 = st.columns(2)
        with c1:
            _cbox("Headline",      ai.get("headline",""))
            _cbox("Subheadline",   ai.get("subheadline",""))
            _cbox("Package Name",  ai.get("package_name",""))
            _cbox("Price",         ai.get("price",""))
            _cbox("Duration",      ai.get("duration",""))
            _cbox("CTA",           ai.get("cta",""))
            _cbox("Highlights",    ai.get("highlights",[]))
            _cbox("Hashtags",      ai.get("hashtags",""))
        with c2:
            _cbox("📸 Instagram Caption",  ai.get("instagram_caption",""))
            _cbox("👥 Facebook Caption",   ai.get("facebook_caption",""))
            _cbox("📱 WhatsApp Status",    ai.get("whatsapp_status",""))
            _cbox("📺 YouTube Title",      ai.get("youtube_title",""))
            _cbox("🎬 Reel Script",        ai.get("reel_script",""))

    # ════════════════════════════════════════════════════════════════════
    # TAB 4 — BULK
    # ════════════════════════════════════════════════════════════════════
    with tab_bulk:
        st.markdown("### 📦 Bulk — All Platforms in One ZIP")
        b_theme_b = st.selectbox("Theme",list(THEMES.keys()),
            index=list(THEMES.keys()).index(ai.get("theme","🏅 Classic Navy Gold"))
                  if ai.get("theme") in THEMES else 0, key="bb_theme")
        b_plats  = st.multiselect("Platforms",list(PLATFORMS.keys()),
                                   default=list(PLATFORMS.keys())[:4])
        b_gen    = st.button("📦 Generate ZIP",type="primary",
                              use_container_width=True,disabled=not b_plats)
        if b_gen and b_plats:
            ai["website"] = ai.get("website") or brand.get("website","")
            ai["phone"]   = ai.get("phone")   or brand.get("phone","")
            ai["_brand"]  = brand
            photos_pil = load_photos(raw)
            qr = _fetch_qr_cached(ai.get("website",""))
            zbuf = io.BytesIO()
            prog = st.progress(0,"Generating banners…")
            preview_done = False
            with zipfile.ZipFile(zbuf,"w",zipfile.ZIP_DEFLATED) as zf:
                for i,pname in enumerate(b_plats):
                    W,H = PLATFORMS[pname]
                    banner = render_banner(photos_pil,W,H,b_theme_b,ai,brand,logo_bytes,qr)
                    safe = re.sub(r"[^\w]","_",pname)[:30]
                    zf.writestr(f"{ai.get('package_name','banner')}_{safe}_{W}x{H}.png",
                                to_bytes(banner,"PNG"))
                    if not preview_done:
                        st.image(to_bytes(banner),caption=pname,use_container_width=True)
                        preview_done = True
                    prog.progress((i+1)/len(b_plats))
            prog.empty(); zbuf.seek(0)
            st.success(f"✅ {len(b_plats)} banners ready!")
            st.download_button("📥 Download ZIP",data=zbuf.getvalue(),
                file_name=f"{ai.get('package_name','banners')}_all_platforms.zip",
                mime="application/zip",use_container_width=True)
        elif not b_plats:
            st.markdown('<div class="empty"><div class="ei">📦</div>'
                        '<div class="et">Select platforms above → Generate</div></div>',
                        unsafe_allow_html=True)

    st.markdown("---")
    with st.expander("💡 Tips + Design guide"):
        st.markdown("""
**One-line prompt tips:**
> *5N/6D Bhutan, Thimphu 2N + Punakha 1N + Paro 2N, Rs 28777/person twin sharing, breakfast dinner, Toyota vehicle, guide, min 4 guests, valid Sept 2026*

Include: destination + city-wise nights + exact price + inclusions + validity

**Theme guide:**
| Package Type | Best Theme |
|---|---|
| Rajasthan / Heritage | 🏅 Classic Navy Gold |
| Wildlife / Nature | 🌿 Emerald Luxury |
| Beach / Andaman / Goa | 🌊 Ocean Blue |
| Luxury / Premium | 🌙 Midnight Premium |
| Adventure / Summer | 🌺 Coral Vibrant |

**Photo tips (for MMT-quality output):**
- Upload 2-4 photos: 1 wide landscape hero + 1-2 monuments/activities
- Min 1500px wide for crisp 2× rendering
- Golden hour / blue hour shots = most dramatic output

**PPTX → Canva workflow:**
Canva → Create Design → Import → Upload PPTX → each element fully editable
        """)
