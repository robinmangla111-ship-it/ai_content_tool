"""
🏖️ Slider Content Creator PRO
Pipeline: Idea → AI Content → Tailwind HTML → Playwright Screenshot → Video → Metadata
Supports English, Hindi, Hinglish
"""

# ── std-lib ────────────────────────────────────────────────────────────────
import os
import io
import re
import json
import time
import random
import zipfile
import tempfile
import base64
from pathlib import Path
from typing import Optional, List, Dict, Tuple

# ── set env BEFORE playwright import ──────────────────────────────────────
# "0" means "use default browser install path", not a custom one.
# Use empty-string removal so system PATH is used.
os.environ.pop("PLAYWRIGHT_BROWSERS_PATH", None)

# ── third-party ───────────────────────────────────────────────────────────
import streamlit as st
import yaml
import numpy as np
from PIL import Image, ImageFilter, ImageEnhance

from pptx import Presentation
from pptx.util import Inches

# ── optional deps ─────────────────────────────────────────────────────────
try:
    from groq import Groq
    GROQ_OK = True
except ImportError:
    GROQ_OK = False

try:
    from huggingface_hub import InferenceClient
    HF_OK = True
except ImportError:
    HF_OK = False

try:
    from moviepy.editor import (
        ImageClip, concatenate_videoclips, vfx, AudioFileClip,
    )
    MOVIEPY_OK = True
except ImportError:
    MOVIEPY_OK = False

try:
    from playwright.sync_api import sync_playwright
    PLAYWRIGHT_OK = True
except ImportError:
    PLAYWRIGHT_OK = False


# ─────────────────────────────────────────────────────────────────────────────
# PAGE CONFIG  (must be first Streamlit call)
# ─────────────────────────────────────────────────────────────────────────────
st.set_page_config(
    page_title="🏖️ Slider Creator PRO",
    page_icon="🏖️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ─────────────────────────────────────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────────────────────────────────────
CONFIG_PATH = Path(__file__).resolve().parent / "config.yaml"

@st.cache_resource
def load_config() -> Dict:
    if not CONFIG_PATH.exists():
        return {}
    with open(CONFIG_PATH, "r", encoding="utf-8") as f:
        return yaml.safe_load(f) or {}

CFG = load_config()

CONTENT_CATEGORIES = [
    "Travel", "Health", "Devotion", "Motivation",
    "Finance", "Education", "General",
]

def _env(name: str) -> str:
    return os.getenv(name, "").strip()


# ─────────────────────────────────────────────────────────────────────────────
# SESSION STATE
# ─────────────────────────────────────────────────────────────────────────────
def init_state() -> None:
    defaults: Dict = {
        "step": 1,
        "topic": "",
        "language": "English",
        "content_category": "Travel",
        "num_slides": 5,
        "slide_content": None,
        "youtube_meta": None,
        "groq_api_key": _env("GROQ_API_KEY"),
        "hf_token": _env("HF_TOKEN"),
        "selected_groq_model": None,
        "png_paths": [],
        "video_path": None,
    }
    for k, v in defaults.items():
        if k not in st.session_state:
            st.session_state[k] = v


# ─────────────────────────────────────────────────────────────────────────────
# CSS
# ─────────────────────────────────────────────────────────────────────────────
def inject_css() -> None:
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Montserrat:wght@400;600;700;800&family=Noto+Sans:wght@400;700;800&display=swap');

    html, body, [class*="css"] {
        font-family: 'Montserrat', 'Noto Sans', sans-serif;
    }
    .main-header {
        background: linear-gradient(135deg,#111827 0%,#1f2937 50%,#0b1220 100%);
        padding: 1.8rem 2rem;
        border-radius: 18px;
        margin-bottom: 1.5rem;
        border: 1px solid rgba(255,165,0,.25);
        box-shadow: 0 8px 30px rgba(0,0,0,.35);
    }
    .main-header h1 { color:#fbbf24; font-size:2rem; margin:0; font-weight:800; letter-spacing:.5px; }
    .main-header p  { color:rgba(255,255,255,.7); margin:.4rem 0 0; font-size:1rem; }
    .step-badge {
        display:inline-block;
        background:linear-gradient(90deg,#fbbf24,#fb7185);
        color:#000; font-weight:800; font-size:.75rem;
        padding:4px 14px; border-radius:20px;
        margin-bottom:.6rem; text-transform:uppercase; letter-spacing:1px;
    }
    .progress-bar  { display:flex; gap:6px; margin-bottom:1.5rem; }
    .progress-step { flex:1; height:6px; border-radius:3px; background:rgba(255,255,255,.1); }
    .progress-step.done   { background:#fbbf24; }
    .progress-step.active { background:rgba(251,191,36,.55); }
    footer { visibility:hidden; }
    </style>
    """, unsafe_allow_html=True)


def progress_bar(current: int, total: int = 5) -> None:
    bars = "".join(
        f'<div class="progress-step {"done" if i < current else "active" if i == current else ""}"></div>'
        for i in range(1, total + 1)
    )
    st.markdown(f'<div class="progress-bar">{bars}</div>', unsafe_allow_html=True)


def step_badge(n: int, label: str) -> None:
    st.markdown(
        f'<div class="step-badge">Step {n}</div>'
        f'<h3 style="color:#fbbf24;margin-top:4px">{label}</h3>',
        unsafe_allow_html=True,
    )


# ─────────────────────────────────────────────────────────────────────────────
# JSON SAFE PARSER
# ─────────────────────────────────────────────────────────────────────────────
def _safe_json(raw: str) -> Dict:
    raw = raw.strip()
    raw = re.sub(r"^```json\s*", "", raw)
    raw = re.sub(r"^```", "", raw)
    raw = re.sub(r"```$", "", raw).strip()
    m = re.search(r"\{.*\}", raw, re.DOTALL)
    if m:
        raw = m.group()
    # strip control chars; fix curly quotes
    raw = re.sub(r"[\x00-\x1F\x7F]", " ", raw)
    raw = raw.replace("\u201c", '"').replace("\u201d", '"').replace("\u2019", "'")
    return json.loads(raw)


# ─────────────────────────────────────────────────────────────────────────────
# CONTENT GENERATION
# ─────────────────────────────────────────────────────────────────────────────
LANG_RULES = {
    "English":  "Generate ALL content strictly in English.",
    "Hindi":    "सभी टेक्स्ट केवल हिंदी (देवनागरी) में लिखें। अंग्रेजी बिल्कुल नहीं।",
    "Hinglish": "Write in Hinglish (Roman Hindi + English mix). Do NOT use Devanagari.",
}

CATEGORY_RULES = {
    "Travel":     "Make it like a professional travel ad. Include packages, highlights, CTA, luxury tone.",
    "Health":     "Educational and safe. Add do/don't. Short disclaimer.",
    "Devotion":   "Spiritual and devotional. Bhakti tone.",
    "Finance":    "Actionable and clear. Add 'Not financial advice' disclaimer.",
    "Motivation": "Punchy, viral, reels-style.",
    "Education":  "Structured, simple to learn.",
    "General":    "Engaging and useful.",
}

DESIGN_TEMPLATES = [
    "luxury_offer", "itinerary", "minimal_quote", "bold_typography",
    "split_modern", "price_focus", "comparison", "checklist_card",
    "devotional_glow", "health_tip",
]


def generate_slide_content(
    topic: str, language: str, num_slides: int,
    api_key: str, model: str, category: str
) -> Optional[Dict]:
    if not GROQ_OK:
        st.error("groq library not installed.  `pip install groq`")
        return None

    prompt = f"""
{LANG_RULES.get(language, LANG_RULES["English"])}

You are a premium YouTube Shorts + Instagram Reels creative director.

CATEGORY : {category}
STYLE    : {CATEGORY_RULES.get(category, CATEGORY_RULES["General"])}
TOPIC    : "{topic}"

Generate exactly {num_slides} slides.

STRICT RULES:
- ALL text fields MUST be in {language}.
- Hindi → Devanagari only.  Hinglish → Roman script only.
- Do NOT put raw newlines inside JSON string values — use \\n instead.
- image_prompt MUST always be in English, cinematic style.

Return ONLY valid JSON — no markdown, no explanation:

{{
  "title": "...",
  "description_short": "...",
  "category": "{category}",
  "slides": [
    {{
      "slide_number": 1,
      "heading":    "...",
      "subheading": "...",
      "content":    "Line1\\nLine2\\nLine3",
      "highlight":  "Badge / offer hook",
      "cta":        "BOOK NOW / SAVE THIS / FOLLOW",
      "price":      "₹19,999 or empty string",
      "speaker_notes": "Narrator text\\nSecond line",
      "image_prompt": "English cinematic photo prompt, 8k, photorealistic",
      "design": {{
        "template": "{random.choice(DESIGN_TEMPLATES)}",
        "accent":   "gold | pink | blue | green | red",
        "text_align": "left | center"
      }}
    }}
  ]
}}
""".strip()

    try:
        client = Groq(api_key=api_key)
        resp = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.75,
            max_tokens=4000,
        )
        return _safe_json(resp.choices[0].message.content)
    except Exception as e:
        st.error(f"Content generation failed: {e}")
        return None


def generate_youtube_meta(
    slide_content: Dict, language: str, category: str,
    api_key: str, model: str,
) -> Optional[Dict]:
    if not GROQ_OK:
        return None

    summary = "\n".join(
        f"Slide {s['slide_number']}: {s['heading']}"
        for s in slide_content.get("slides", [])
    )

    prompt = f"""
{LANG_RULES.get(language, LANG_RULES["English"])}

SEO metadata for YouTube Shorts + Instagram Reels.

Slides:
{summary}

Return ONLY JSON (no markdown):
{{
  "title": "...",
  "description": "...",
  "tags": ["..."],
  "hashtags": ["#..."],
  "thumbnail_text": "...",
  "thumbnail_subtext": "...",
  "seo_keywords": ["..."]
}}

Minimum: 20 tags, 12 hashtags.
""".strip()

    try:
        client = Groq(api_key=api_key)
        resp = client.chat.completions.create(
            model=model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=2000,
        )
        return _safe_json(resp.choices[0].message.content)
    except Exception as e:
        st.error(f"Meta generation failed: {e}")
        return None


# ─────────────────────────────────────────────────────────────────────────────
# HF IMAGE GENERATION (optional helper)
# ─────────────────────────────────────────────────────────────────────────────
def generate_image_hf(
    prompt: str, hf_token: str, model_id: str
) -> Optional[Image.Image]:
    if not (HF_OK and hf_token):
        return None
    try:
        client = InferenceClient(token=hf_token)
        img = client.text_to_image(prompt, model=model_id)
        return img if isinstance(img, Image.Image) else Image.open(io.BytesIO(img))
    except Exception as e:
        st.warning(f"HF image gen failed: {e}")
        return None


# ─────────────────────────────────────────────────────────────────────────────
# IMAGE → BASE64  (used for Playwright embedding — avoids file:// CSP issues)
# ─────────────────────────────────────────────────────────────────────────────
def img_to_data_url(path: str) -> str:
    """Convert a local image to a data URL so Playwright can embed it inline."""
    ext = Path(path).suffix.lower().lstrip(".")
    mime = {"jpg": "image/jpeg", "jpeg": "image/jpeg",
            "png": "image/png",  "webp": "image/webp"}.get(ext, "image/png")
    data = Path(path).read_bytes()
    b64  = base64.b64encode(data).decode()
    return f"data:{mime};base64,{b64}"


# ─────────────────────────────────────────────────────────────────────────────
# HTML TEMPLATE ENGINE  (Tailwind 1080 × 1920 — portrait / Shorts)
# ─────────────────────────────────────────────────────────────────────────────
TAILWIND_CDN = "https://cdn.tailwindcss.com"


def _accent_colors(accent: str) -> Tuple[str, str]:
    palette = {
        "pink":  ("#fb7185", "#f472b6"),
        "blue":  ("#38bdf8", "#6366f1"),
        "green": ("#22c55e", "#10b981"),
        "red":   ("#fb7185", "#ef4444"),
    }
    return palette.get((accent or "").lower().strip(), ("#fbbf24", "#f97316"))


def build_slide_html(slide: Dict, bg_img_path: Optional[str] = None) -> str:
    heading    = slide.get("heading", "")
    subheading = slide.get("subheading", "")
    raw_content = slide.get("content", "")
    # content may use literal \n (from JSON) or actual newlines
    content    = [l.strip() for l in re.split(r"\\n|\n", raw_content) if l.strip()]
    highlight  = slide.get("highlight", "")
    cta        = slide.get("cta", "")
    price      = slide.get("price", "")

    design    = slide.get("design") or {}
    template  = design.get("template", "")
    accent    = design.get("accent", "gold")
    text_align = design.get("text_align", "left")

    if template not in DESIGN_TEMPLATES:
        template = random.choice(DESIGN_TEMPLATES)

    c1, c2 = _accent_colors(accent)

    # ── background style ──────────────────────────────────────────────────
    # Use data URL to avoid Playwright file:// CSP restrictions
    if bg_img_path and os.path.exists(bg_img_path):
        data_url = img_to_data_url(bg_img_path)
        bg_css = (
            f"background-image:url('{data_url}');"
            "background-size:cover;background-position:center;"
        )
    else:
        bg_css = "background-color:#0b1220;"

    align_class = "text-center items-center" if text_align == "center" else "text-left items-start"

    def bullets(lines: List[str], icon: str = "✔") -> str:
        return "".join(
            f'<div class="flex gap-4 items-start">'
            f'  <div class="text-[48px] font-extrabold shrink-0" style="color:{c1}">{icon}</div>'
            f'  <div class="text-[46px] font-semibold text-white/90 leading-snug">{ln}</div>'
            f'</div>'
            for ln in lines[:6] if ln
        )

    # ── template body blocks ──────────────────────────────────────────────
    if template == "luxury_offer":
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-b from-black/60 via-black/40 to-black/75"></div>
        <div class="relative px-16 pt-20 flex flex-col justify-between h-full">
          <div>
            <div class="inline-flex items-center gap-3 px-10 py-4 rounded-full font-extrabold text-black text-[38px] shadow-2xl"
                 style="background:linear-gradient(90deg,{c1},{c2})">✨ {highlight or "LIMITED OFFER"}</div>
            <div class="mt-10 text-[100px] font-extrabold leading-[1.02] tracking-tight drop-shadow-2xl">{heading}</div>
            <div class="mt-6 text-[50px] font-semibold text-white/80 max-w-[900px]">{subheading}</div>
          </div>
          <div class="mb-20 bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[60px] p-14 shadow-2xl">
            <div class="space-y-8">{bullets(content[:5])}</div>
            <div class="mt-12 flex justify-between items-center">
              <div class="text-[68px] font-extrabold" style="color:{c1}">{price}</div>
              <div class="px-14 py-5 rounded-[50px] text-[46px] font-extrabold text-black shadow-2xl"
                   style="background:linear-gradient(90deg,{c1},{c2})">{cta or "BOOK NOW"}</div>
            </div>
          </div>
        </div>"""

    elif template == "itinerary":
        items = "".join(
            f'<div class="flex gap-5 items-start">'
            f'  <div class="w-[72px] h-[72px] rounded-2xl flex items-center justify-center font-extrabold text-[38px] text-black shadow-xl shrink-0"'
            f'       style="background:linear-gradient(90deg,{c1},{c2})">{idx}</div>'
            f'  <div class="text-[46px] font-semibold text-white/90 leading-snug">{ln}</div>'
            f'</div>'
            for idx, ln in enumerate(content[:6], 1)
        )
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-b from-black/65 via-black/40 to-black/80"></div>
        <div class="relative px-16 pt-20">
          <div class="flex justify-between items-start">
            <div>
              <div class="text-[92px] font-extrabold leading-[1.05] drop-shadow-xl">{heading}</div>
              <div class="mt-6 text-[46px] text-white/80 font-semibold">{subheading}</div>
            </div>
            <div class="px-10 py-4 rounded-full text-[40px] font-extrabold text-black shadow-2xl"
                 style="background:linear-gradient(90deg,{c1},{c2})">{highlight or "ITINERARY"}</div>
          </div>
          <div class="mt-20 bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[55px] p-14 space-y-10 shadow-2xl">{items}</div>
          <div class="mt-14 flex justify-between items-center">
            <div class="text-[64px] font-extrabold" style="color:{c1}">{price}</div>
            <div class="px-14 py-5 rounded-[50px] text-[46px] font-extrabold text-black shadow-2xl"
                 style="background:linear-gradient(90deg,{c1},{c2})">{cta or "BOOK NOW"}</div>
          </div>
        </div>"""

    elif template == "minimal_quote":
        body = f"""
        <div class="absolute inset-0 bg-black/70"></div>
        <div class="relative flex flex-col justify-center h-full px-20 {align_class}">
          <div class="text-[110px] font-extrabold leading-[1.05] tracking-tight drop-shadow-2xl">{heading}</div>
          <div class="mt-10 text-white/85 text-[50px] font-semibold max-w-[900px] leading-snug">{" ".join(content[:3])}</div>
          <div class="mt-16 px-14 py-5 rounded-full font-extrabold text-black text-[44px] shadow-2xl"
               style="background:linear-gradient(90deg,{c1},{c2})">{cta or "FOLLOW FOR MORE"}</div>
        </div>"""

    elif template == "bold_typography":
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-tr from-black/80 via-black/40 to-black/90"></div>
        <div class="relative flex flex-col justify-between h-full px-16 pt-20 pb-20">
          <div>
            <div class="text-[130px] font-extrabold leading-[0.95] tracking-tight drop-shadow-2xl">{heading}</div>
            <div class="mt-8 text-[52px] font-semibold text-white/80 max-w-[950px]">{subheading}</div>
          </div>
          <div class="bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[55px] p-14 shadow-2xl space-y-8">{bullets(content[:4], "➤")}</div>
          <div class="flex justify-between items-center mt-12">
            <div class="px-10 py-4 rounded-full font-extrabold text-black text-[40px] shadow-2xl"
                 style="background:linear-gradient(90deg,{c1},{c2})">{highlight or "TRENDING"}</div>
            <div class="text-[62px] font-extrabold" style="color:{c1}">{price}</div>
          </div>
        </div>"""

    elif template == "split_modern":
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-b from-black/30 via-black/60 to-black/85"></div>
        <div class="relative grid grid-cols-2 h-full">
          <div class="p-16 flex flex-col justify-between">
            <div>
              <div class="inline-block px-10 py-4 rounded-full font-extrabold text-black text-[38px] shadow-2xl"
                   style="background:linear-gradient(90deg,{c1},{c2})">{highlight or "SPECIAL DEAL"}</div>
              <div class="mt-10 text-[88px] font-extrabold leading-[1.05] drop-shadow-2xl">{heading}</div>
              <div class="mt-6 text-[48px] font-semibold text-white/80">{subheading}</div>
            </div>
            <div class="text-[66px] font-extrabold" style="color:{c1}">{price}</div>
          </div>
          <div class="p-16 flex flex-col justify-end">
            <div class="bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[60px] p-14 shadow-2xl space-y-8">{bullets(content[:5])}</div>
            <div class="mt-10 px-14 py-5 rounded-[50px] text-[46px] font-extrabold text-black shadow-2xl text-center"
                 style="background:linear-gradient(90deg,{c1},{c2})">{cta or "BOOK NOW"}</div>
          </div>
        </div>"""

    elif template == "price_focus":
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-b from-black/75 via-black/40 to-black/90"></div>
        <div class="relative flex flex-col justify-center h-full px-16 text-center items-center">
          <div class="px-12 py-5 rounded-full font-extrabold text-black text-[42px] shadow-2xl"
               style="background:linear-gradient(90deg,{c1},{c2})">{highlight or "FLASH SALE"}</div>
          <div class="mt-12 text-[110px] font-extrabold leading-[1.0] drop-shadow-2xl">{heading}</div>
          <div class="mt-6 text-[56px] font-semibold text-white/80 max-w-[950px]">{subheading}</div>
          <div class="mt-16 text-[140px] font-extrabold tracking-tight drop-shadow-2xl" style="color:{c1}">{price}</div>
          <div class="mt-14 bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[60px] p-14 shadow-2xl w-full max-w-[950px] space-y-8 text-left">{bullets(content[:4])}</div>
          <div class="mt-16 px-16 py-6 rounded-[60px] text-[52px] font-extrabold text-black shadow-2xl"
               style="background:linear-gradient(90deg,{c1},{c2})">{cta or "BOOK NOW"}</div>
        </div>"""

    elif template == "comparison":
        left, right = content[:3], content[3:6]

        def cmp_block(title: str, lines: List[str], color: str) -> str:
            items = "".join(f'<div class="text-[42px] font-semibold text-white/90">• {l}</div>' for l in lines)
            return (
                f'<div class="bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[55px] p-12 shadow-2xl">'
                f'  <div class="text-[54px] font-extrabold mb-6" style="color:{color}">{title}</div>'
                f'  <div class="space-y-5">{items}</div>'
                f'</div>'
            )

        body = f"""
        <div class="absolute inset-0 bg-gradient-to-br from-black/80 via-black/40 to-black/90"></div>
        <div class="relative px-16 pt-20">
          <div class="text-[92px] font-extrabold drop-shadow-2xl">{heading}</div>
          <div class="mt-6 text-[48px] font-semibold text-white/80">{subheading}</div>
          <div class="mt-14 grid grid-cols-2 gap-10">
            {cmp_block("Option A", left, c1)}
            {cmp_block("Option B", right, c2)}
          </div>
          <div class="mt-12 flex justify-between items-center">
            <div class="text-[66px] font-extrabold" style="color:{c1}">{price}</div>
            <div class="px-14 py-5 rounded-[55px] text-[46px] font-extrabold text-black shadow-2xl"
                 style="background:linear-gradient(90deg,{c1},{c2})">{cta or "CHOOSE NOW"}</div>
          </div>
        </div>"""

    elif template == "checklist_card":
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-b from-black/65 via-black/40 to-black/90"></div>
        <div class="relative px-16 pt-20">
          <div class="inline-block px-10 py-4 rounded-full font-extrabold text-black text-[40px] shadow-2xl"
               style="background:linear-gradient(90deg,{c1},{c2})">{highlight or "MUST KNOW"}</div>
          <div class="mt-10 text-[100px] font-extrabold drop-shadow-2xl leading-[1.02]">{heading}</div>
          <div class="mt-6 text-[52px] font-semibold text-white/80 max-w-[950px]">{subheading}</div>
          <div class="mt-14 bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[65px] p-16 shadow-2xl space-y-10">{bullets(content[:6], "☑")}</div>
          <div class="mt-14 px-16 py-6 rounded-[60px] text-[50px] font-extrabold text-black shadow-2xl text-center"
               style="background:linear-gradient(90deg,{c1},{c2})">{cta or "SAVE THIS"}</div>
        </div>"""

    elif template == "devotional_glow":
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-b from-black/80 via-black/55 to-black/90"></div>
        <div class="absolute top-[-200px] left-[-200px] w-[600px] h-[600px] rounded-full blur-[120px]" style="background:{c1};opacity:.35;"></div>
        <div class="absolute bottom-[-250px] right-[-250px] w-[700px] h-[700px] rounded-full blur-[150px]" style="background:{c2};opacity:.25;"></div>
        <div class="relative flex flex-col justify-center h-full px-16 text-center items-center">
          <div class="text-[100px] font-extrabold drop-shadow-2xl">{heading}</div>
          <div class="mt-8 text-[54px] font-semibold text-white/80 max-w-[950px]">{subheading}</div>
          <div class="mt-14 bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[65px] p-16 shadow-2xl w-full max-w-[950px]">
            <div class="text-[50px] font-semibold text-white/90 leading-snug">{" ".join(content[:4])}</div>
          </div>
          <div class="mt-16 px-16 py-6 rounded-[60px] text-[52px] font-extrabold text-black shadow-2xl"
               style="background:linear-gradient(90deg,{c1},{c2})">{cta or "🙏 जय श्री राम"}</div>
        </div>"""

    else:  # health_tip (default fallback)
        body = f"""
        <div class="absolute inset-0 bg-gradient-to-b from-black/75 via-black/45 to-black/90"></div>
        <div class="relative px-16 pt-20">
          <div class="inline-block px-10 py-4 rounded-full font-extrabold text-black text-[40px] shadow-2xl"
               style="background:linear-gradient(90deg,{c1},{c2})">{highlight or "HEALTH TIP"}</div>
          <div class="mt-10 text-[100px] font-extrabold drop-shadow-2xl leading-[1.02]">{heading}</div>
          <div class="mt-6 text-[52px] font-semibold text-white/80 max-w-[950px]">{subheading}</div>
          <div class="mt-14 bg-white/10 backdrop-blur-2xl border border-white/15 rounded-[65px] p-16 shadow-2xl space-y-10">{bullets(content[:6])}</div>
          <div class="mt-14 px-16 py-6 rounded-[60px] text-[50px] font-extrabold text-black shadow-2xl text-center"
               style="background:linear-gradient(90deg,{c1},{c2})">{cta or "SAVE THIS"}</div>
        </div>"""

    return f"""<!DOCTYPE html>
<html>
<head>
  <meta charset="utf-8">
  <script src="{TAILWIND_CDN}"></script>
  <style>
    body {{
      margin:0; width:1080px; height:1920px;
      overflow:hidden;
      font-family:'Montserrat','Noto Sans',sans-serif;
      color:white;
    }}
    .bg {{
      position:absolute; inset:0;
      {bg_css}
      filter:contrast(1.1) saturate(1.1);
    }}
    .vignette {{
      position:absolute; inset:0;
      background:radial-gradient(circle at center,rgba(0,0,0,.1),rgba(0,0,0,.75));
    }}
  </style>
</head>
<body>
  <div class="bg"></div>
  <div class="vignette"></div>
  {body}
  <div class="absolute bottom-8 left-14 text-white/40 text-[34px] font-semibold">
    Premium Shorts • Slide {slide.get('slide_number','')}
  </div>
</body>
</html>"""


# ─────────────────────────────────────────────────────────────────────────────
# PLAYWRIGHT RENDER
# ─────────────────────────────────────────────────────────────────────────────
def render_html_to_png(html: str, out_path: str) -> bool:
    if not PLAYWRIGHT_OK:
        st.error("Playwright not installed.  Run: `pip install playwright && playwright install chromium`")
        return False
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(
                headless=True,
                args=["--no-sandbox", "--disable-dev-shm-usage", "--disable-web-security"],
            )
            page = browser.new_page(viewport={"width": 1080, "height": 1920})
            page.set_content(html, wait_until="networkidle")
            page.wait_for_timeout(2000)   # let Tailwind JIT render
            page.screenshot(path=out_path, clip={"x": 0, "y": 0, "width": 1080, "height": 1920})
            browser.close()
        return True
    except Exception as e:
        st.error(f"Playwright render failed: {e}")
        return False


# ─────────────────────────────────────────────────────────────────────────────
# PPTX
# ─────────────────────────────────────────────────────────────────────────────
def create_pptx_from_pngs(png_paths: List[str]) -> bytes:
    prs = Presentation()
    prs.slide_width  = Inches(9)
    prs.slide_height = Inches(16)
    for path in png_paths:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(path, 0, 0, prs.slide_width, prs.slide_height)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────────────────────
# VIDEO  (Ken Burns + cross-dissolve — simple & robust)
# ─────────────────────────────────────────────────────────────────────────────
def _ken_burns(clip: "ImageClip", zoom: float = 1.08, direction: str = "in") -> "ImageClip":
    """
    Correct MoviePy fl() usage: inner function signature is (get_frame, t).
    Returns a clip with a gentle zoom-pan effect.
    """
    w, h = clip.size
    dur  = clip.duration

    def make_frame(get_frame, t: float):          # ← FIXED: two args
        frame = get_frame(t)
        img   = Image.fromarray(frame.astype(np.uint8))

        progress = t / dur if dur > 0 else 0
        z = (1 + (zoom - 1) * progress) if direction == "in" else (zoom - (zoom - 1) * progress)
        z = max(z, 1.0)

        new_w, new_h = int(w * z), int(h * z)
        img = img.resize((new_w, new_h), Image.LANCZOS)

        # centre-crop back to original size
        x0 = (new_w - w) // 2
        y0 = (new_h - h) // 2
        img = img.crop((x0, y0, x0 + w, y0 + h))
        return np.array(img)

    return clip.fl(make_frame)                    # apply_to not needed


def build_video_from_pngs(
    png_paths: List[str],
    duration: int,
    output_path: str,
    bg_music_path: Optional[str] = None,
    music_volume: float = 0.15,
) -> bool:
    if not MOVIEPY_OK:
        return False

    clips = []
    for path in png_paths:
        base = ImageClip(path).set_duration(duration)
        direction = random.choice(["in", "out"])
        zoom      = random.choice([1.05, 1.08, 1.12])
        animated  = _ken_burns(base, zoom=zoom, direction=direction)
        # simple fade in/out  (no chained swipe — avoids the closure bug)
        animated  = animated.fx(vfx.fadein, 0.3).fx(vfx.fadeout, 0.3)
        clips.append(animated)

    if not clips:
        return False

    # Cross-dissolve concatenation using crossfadein (MoviePy built-in)
    XFADE = 0.4
    for i in range(1, len(clips)):
        clips[i] = clips[i].crossfadein(XFADE)

    final = concatenate_videoclips(clips, padding=-XFADE, method="compose")

    # Background music
    if bg_music_path and os.path.exists(bg_music_path):
        try:
            music = AudioFileClip(bg_music_path).volumex(music_volume)
            if music.duration < final.duration:
                reps  = int(final.duration // music.duration) + 1
                music = concatenate_videoclips([music] * reps).audio   # audio loop
            music = music.subclip(0, final.duration)
            final = final.set_audio(music)
        except Exception as e:
            st.warning(f"Music attachment failed (video still saved): {e}")

    final.write_videofile(
        output_path,
        fps=30,
        codec="libx264",
        audio_codec="aac",
        bitrate="9000k",
        threads=4,
        logger=None,
    )
    final.close()
    return True


# ─────────────────────────────────────────────────────────────────────────────
# SIDEBAR
# ─────────────────────────────────────────────────────────────────────────────
def sidebar() -> None:
    with st.sidebar:
        st.markdown("## 🏖️ Slider Creator PRO")
        st.markdown("---")
        st.markdown("### 📦 Library Status")
        rows = [
            ("Groq",       GROQ_OK),
            ("HuggingFace",HF_OK),
            ("Playwright", PLAYWRIGHT_OK),
            ("MoviePy",    MOVIEPY_OK),
        ]
        for name, ok in rows:
            icon = "✅" if ok else "❌"
            st.write(f"{icon}  {name}")

        st.markdown("---")
        steps = {
            1: "📝 Topic & Settings",
            2: "✏️ Review Content",
            3: "🎨 Render Slides",
            4: "📥 Export Assets",
            5: "📊 SEO Metadata",
        }
        st.markdown("### 🗺️ Navigation")
        for n, label in steps.items():
            if st.button(label, key=f"nav_s{n}", use_container_width=True,
                         disabled=(n > 1 and not st.session_state.slide_content)):
                st.session_state.step = n
                st.rerun()

        st.markdown("---")
        if not PLAYWRIGHT_OK:
            st.warning("Install Playwright:\n```\npip install playwright\nplaywright install chromium\n```")

        if st.button("🔄 Reset App", use_container_width=True):
            for k in list(st.session_state.keys()):
                del st.session_state[k]
            st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 1 — Topic & Settings
# ─────────────────────────────────────────────────────────────────────────────
def page_step1() -> None:
    step_badge(1, "Topic & Settings")
    progress_bar(1)

    c1, c2 = st.columns([3, 2])
    with c1:
        st.session_state.topic = st.text_area(
            "🎯 Slider topic / idea",
            value=st.session_state.topic,
            height=120,
            placeholder="Example: Dubai Travel Package Offer",
        )
        st.session_state.language = st.radio(
            "🌐 Language",
            ["English", "Hindi", "Hinglish"],
            index=["English", "Hindi", "Hinglish"].index(st.session_state.language),
            horizontal=True,
        )
        st.session_state.content_category = st.selectbox(
            "📌 Content Type", CONTENT_CATEGORIES,
            index=CONTENT_CATEGORIES.index(st.session_state.content_category),
        )

    with c2:
        st.session_state.num_slides = st.slider("Slides", 3, 10, st.session_state.num_slides)
        models = CFG.get("groq", {}).get("models", ["llama3-8b-8192"])
        model = st.selectbox("🤖 Groq Model", models, index=0)
        st.session_state.selected_groq_model = model

        st.markdown("---")
        st.session_state.groq_api_key = st.text_input(
            "Groq API Key *", type="password", value=st.session_state.groq_api_key,
            placeholder="gsk_…",
        )
        st.session_state.hf_token = st.text_input(
            "HF Token (optional)", type="password", value=st.session_state.hf_token,
            placeholder="hf_…",
        )

    ready = bool(st.session_state.topic.strip() and st.session_state.groq_api_key.strip())
    if not ready:
        st.info("💡 Enter your topic and Groq API key to continue.")

    if st.button("🚀 Generate Slide Content", disabled=not ready):
        with st.spinner("Generating premium content via Groq…"):
            result = generate_slide_content(
                st.session_state.topic,
                st.session_state.language,
                st.session_state.num_slides,
                st.session_state.groq_api_key,
                st.session_state.selected_groq_model,
                st.session_state.content_category,
            )
        if result:
            st.session_state.slide_content = result
            st.success("✅ Content generated!")
            time.sleep(0.5)
            st.session_state.step = 2
            st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 2 — Review & Edit
# ─────────────────────────────────────────────────────────────────────────────
def page_step2() -> None:
    step_badge(2, "Review, Edit & Upload Images")
    progress_bar(2)

    sc = st.session_state.slide_content
    if not sc:
        st.warning("No content yet.  Go back to Step 1.")
        return

    st.markdown(f"### 🎬 `{sc.get('title','')}`")
    st.caption(sc.get("description_short", ""))

    slides = sc.get("slides", [])
    for i, s in enumerate(slides):
        with st.expander(f"Slide {s.get('slide_number', i+1)} — {s.get('heading','')}", expanded=(i == 0)):
            slides[i]["heading"]    = st.text_input("Heading",    value=s.get("heading", ""),    key=f"h_{i}")
            slides[i]["subheading"] = st.text_input("Subheading", value=s.get("subheading", ""), key=f"sh_{i}")
            slides[i]["content"]    = st.text_area("Content (use \\n for new lines)",
                                                    value=s.get("content", ""), height=130, key=f"c_{i}")
            col1, col2, col3 = st.columns(3)
            slides[i]["highlight"] = col1.text_input("Badge",  value=s.get("highlight", ""), key=f"hl_{i}")
            slides[i]["cta"]       = col2.text_input("CTA",    value=s.get("cta", ""),       key=f"cta_{i}")
            slides[i]["price"]     = col3.text_input("Price",  value=s.get("price", ""),     key=f"pr_{i}")
            slides[i]["image_prompt"] = st.text_input(
                "Image Prompt (English)", value=s.get("image_prompt", ""), key=f"ip_{i}"
            )

            up = st.file_uploader(
                f"📷 Background Image for Slide {i+1} (optional)",
                type=["png", "jpg", "jpeg"], key=f"upload_slide_{i}",
            )
            if up:
                dest = str(Path(tempfile.gettempdir()) / f"user_slide_{i+1:02d}.png")
                Image.open(up).convert("RGB").save(dest)
                slides[i]["uploaded_bg"] = dest
                st.success("✅ Image saved.")

            if slides[i].get("uploaded_bg") and os.path.exists(slides[i]["uploaded_bg"]):
                st.image(slides[i]["uploaded_bg"], caption="Background preview", use_container_width=True)

    sc["slides"] = slides
    st.session_state.slide_content = sc

    if st.button("➡ Next: Render Slides"):
        st.session_state.step = 3
        st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 3 — Render
# ─────────────────────────────────────────────────────────────────────────────
def page_step3() -> None:
    step_badge(3, "Render Pro Slides (Tailwind + Playwright)")
    progress_bar(3)

    if not PLAYWRIGHT_OK:
        st.error(
            "Playwright is not installed or Chromium is missing.\n\n"
            "**Fix:**\n```bash\npip install playwright\nplaywright install chromium\n```"
        )
        return

    sc = st.session_state.slide_content
    if not sc:
        st.warning("No content.  Go back to Step 1.")
        return

    slides = sc.get("slides", [])
    st.info(f"Ready to render **{len(slides)} slides** as 1080×1920 PNG (portrait / Shorts format).")

    # Show previously rendered previews
    if st.session_state.get("png_paths"):
        st.markdown("#### Previously Rendered Slides")
        cols = st.columns(min(len(st.session_state["png_paths"]), 3))
        for j, p in enumerate(st.session_state["png_paths"]):
            with cols[j % 3]:
                st.image(p, use_container_width=True, caption=f"Slide {j+1}")

    if st.button("🎨 Render All Slides → PNG"):
        out_dir = Path(tempfile.gettempdir()) / "slider_html_slides"
        out_dir.mkdir(exist_ok=True)

        # clean stale files
        for old in out_dir.glob("slide_*.png"):
            try: old.unlink()
            except: pass

        prog = st.progress(0.0, text="Starting render…")
        png_paths: List[str] = []
        errors = 0

        for i, slide in enumerate(slides):
            prog.progress((i + 0.1) / len(slides), text=f"Rendering slide {i+1}/{len(slides)}…")
            bg_img = slide.get("uploaded_bg")
            html   = build_slide_html(slide, bg_img_path=bg_img)
            out_png = str(out_dir / f"slide_{i+1:02d}.png")

            if render_html_to_png(html, out_png):
                png_paths.append(out_png)
            else:
                errors += 1

            prog.progress((i + 1) / len(slides))

        st.session_state["png_paths"] = png_paths
        prog.progress(1.0, text="Done!")

        if errors:
            st.warning(f"⚠️ {errors} slide(s) failed to render.")
        if png_paths:
            st.success(f"✅ {len(png_paths)} slide(s) rendered!")
            st.session_state.step = 4
            st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 4 — Export
# ─────────────────────────────────────────────────────────────────────────────
def page_step4() -> None:
    step_badge(4, "Export Assets (PNG / PPTX / MP4)")
    progress_bar(4)

    png_paths: List[str] = [p for p in st.session_state.get("png_paths", []) if os.path.exists(p)]
    if not png_paths:
        st.warning("No rendered slides found.  Go back to Step 3.")
        return

    # ── preview ──────────────────────────────────────────────────────────
    st.markdown("#### Slide Preview")
    cols = st.columns(min(len(png_paths), 4))
    for j, p in enumerate(png_paths):
        with cols[j % 4]:
            st.image(p, use_container_width=True, caption=f"Slide {j+1}")

    st.markdown("---")

    # ── ZIP download ─────────────────────────────────────────────────────
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for p in png_paths:
            zf.write(p, arcname=os.path.basename(p))
    st.download_button("⬇️ Download PNG Slides (ZIP)", zip_buf.getvalue(),
                       "slider_slides.zip", "application/zip")

    # ── PPTX ─────────────────────────────────────────────────────────────
    st.download_button(
        "⬇️ Download PPTX",
        create_pptx_from_pngs(png_paths),
        "slider_slides.pptx",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    # ── Video ─────────────────────────────────────────────────────────────
    st.markdown("---")
    st.markdown("### 🎬 Build MP4 Video")
    if not MOVIEPY_OK:
        st.warning("MoviePy not installed.  `pip install moviepy`  → MP4 export disabled.")
    else:
        duration   = st.slider("Seconds per slide", 2, 10, 5)
        bg_music   = st.file_uploader("🎵 Optional background MP3", type=["mp3"])
        music_path: Optional[str] = None
        if bg_music:
            music_path = str(Path(tempfile.gettempdir()) / "bg_music.mp3")
            with open(music_path, "wb") as f:
                f.write(bg_music.getbuffer())
            st.success("✅ Music file ready.")

        if st.button("🎬 Build Cinematic MP4"):
            out_mp4 = str(Path(tempfile.gettempdir()) / "slider_video.mp4")
            with st.spinner("Encoding video with Ken Burns + cross-dissolve transitions…"):
                ok = build_video_from_pngs(png_paths, duration, out_mp4, bg_music_path=music_path)

            if ok:
                st.session_state.video_path = out_mp4
                st.success("✅ MP4 ready!")
                st.video(out_mp4)
                with open(out_mp4, "rb") as f:
                    st.download_button("⬇️ Download MP4", f.read(), "slider_video.mp4", "video/mp4")
            else:
                st.error("Video build failed.  Check FFmpeg is installed (`ffmpeg -version`).")

    st.markdown("---")
    if st.button("➡ Next: Generate SEO Metadata"):
        st.session_state.step = 5
        st.rerun()


# ─────────────────────────────────────────────────────────────────────────────
# STEP 5 — SEO Metadata
# ─────────────────────────────────────────────────────────────────────────────
def page_step5() -> None:
    step_badge(5, "SEO Metadata — YouTube Shorts / Reels")
    progress_bar(5)

    sc = st.session_state.slide_content
    if not sc:
        st.warning("No slide content found.  Go back to Step 1.")
        return

    if st.button("🤖 Generate SEO Metadata via Groq"):
        with st.spinner("Generating metadata…"):
            meta = generate_youtube_meta(
                sc,
                st.session_state.language,
                st.session_state.content_category,
                st.session_state.groq_api_key,
                st.session_state.selected_groq_model,
            )
        if meta:
            st.session_state.youtube_meta = meta
            st.success("✅ Metadata generated!")

    meta = st.session_state.youtube_meta
    if meta:
        c1, c2 = st.columns([3, 2])
        with c1:
            meta["title"]       = st.text_input("📌 Title",       value=meta.get("title", ""))
            meta["description"] = st.text_area("📄 Description", value=meta.get("description", ""), height=200)
            tags_raw            = st.text_area("🏷️ Tags (comma-separated)",
                                               value=", ".join(meta.get("tags", [])), height=80)
            meta["tags"] = [t.strip() for t in tags_raw.split(",") if t.strip()]

        with c2:
            st.markdown("**# Hashtags**")
            st.code(" ".join(meta.get("hashtags", [])))
            st.markdown("**🖼️ Thumbnail**")
            st.info(f"**{meta.get('thumbnail_text','')}**\n\n{meta.get('thumbnail_subtext','')}")
            st.markdown("**🔍 SEO Keywords**")
            st.write(", ".join(meta.get("seo_keywords", [])))

        st.session_state.youtube_meta = meta
        st.markdown("---")
        export = {"slides": sc.get("slides", []), "metadata": meta}
        st.download_button(
            "⬇️ Download Metadata JSON",
            data=json.dumps(export, ensure_ascii=False, indent=2),
            file_name="slider_metadata.json",
            mime="application/json",
        )


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def render() -> None:
    inject_css()
    init_state()
    sidebar()

    st.markdown("""
    <div class="main-header">
        <h1>🏖️ Slider Content Creator PRO</h1>
        <p>Premium Tailwind + Playwright designs for YouTube Shorts / Instagram Reels &nbsp;|&nbsp;
           English • हिंदी • Hinglish</p>
    </div>
    """, unsafe_allow_html=True)

    step = st.session_state.step
    {1: page_step1, 2: page_step2, 3: page_step3, 4: page_step4, 5: page_step5}.get(step, page_step1)()


if __name__ == "__main__":
    render()
